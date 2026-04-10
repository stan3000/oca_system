import streamlit as st
import os
from pathlib import Path

# =========================================================
# MUST BE FIRST STREAMLIT COMMAND
# =========================================================
st.set_page_config(page_title="OCA Attendance Dashboard", layout="wide")


import streamlit as st

# =========================================================
# HARD-CODED USER ACCOUNTS (NO SECRETS REQUIRED)
# =========================================================
USERS = {
    "admin": "password123",
    "treasurer": "finance2024",
    "secretary": "attend2024"
}

# =========================================================
# LOGIN PAGE
# =========================================================
def login_page():
    st.title("🔐 Login Required")

    username = st.text_input("Username")
    password = st.text_input("Password", type="password")

    if st.button("Login"):
        if username in USERS and USERS[username] == password:
            st.session_state["logged_in"] = True
            st.session_state["user_role"] = username
            st.success("Login successful! Redirecting…")
            st.rerun()
        else:
            st.error("❌ Invalid username or password")


# =========================================================
# MAIN DASHBOARD
# =========================================================
def main_dashboard():
    st.sidebar.markdown(f"👤 Logged in as: **{st.session_state['user_role']}**")

    # logout
    if st.sidebar.button("🚪 Logout"):
        st.session_state["logged_in"] = False
        st.session_state["user_role"] = None
        st.rerun()

    # NAVIGATION
    page = st.sidebar.write("Welcome to OCA System")

    # ==============================
    # HOME PAGE
    # ==============================
    if page == "🏠 Home":
        st.title("🏠 OCA Dashboard")
        st.write("Welcome to the Owerri Cultural Association System!")

    # ==============================
    # ATTENDANCE PAGE
    # ==============================
    elif page == "🧾 Attendance Tracker":
        attendance_ui()  # <-- your existing function remains

    # ==============================
    # FINANCE PAGE
    # ==============================
    elif page == "💰 OCA Finance":
        financial_ui()  # <-- your existing function remains


    # ==============================
    # OCA DUES
    # ==============================

    elif page == "OCA Dues Analysis":
        oca_dues_ui()

# =========================================================
# LOGIN CONTROLLER
# =========================================================
if "logged_in" not in st.session_state:
    st.session_state["logged_in"] = False
    st.session_state["user_role"] = None

if not st.session_state["logged_in"]:
    login_page()
    st.stop()

main_dashboard()



# =========================== LOGIN ENDS HERE ===================
# =========================== LOGIN ENDS HERE ===================
# =========================== LOGIN ENDS HERE ===================





















# # ================================== LOGIN ENDS HERE
#
# import streamlit as st
# import pandas as pd
# from datetime import date
# from io import BytesIO
# import altair as alt
#
# # 🚀 MUST BE FIRST STREAMLIT COMMAND


# # --------------------------------------------------------------------------
# # Sidebar navigation
# # --------------------------------------------------------------------------
# def main():
#     st.sidebar.title("📋 Data Selection")
#     page_selection = st.sidebar.radio(
#         "Choose a page:",
#         ["🏠 Home", "🧾 Attendance Tracker", "🧾 OCA FINANCIAL"]
#     )
#
#     if page_selection == "🏠 Home":
#         st.title("🏠 OCA Dashboard")
#         st.markdown("""
#         Welcome to the **Owerri Cultural Association Attendance Tracker**
#         Use the sidebar to navigate to:
#         - 🧾 Attendance Tracker
#         """)
#
#     elif page_selection == "🧾 Attendance Tracker":
#         attendance_ui()
#
#
#
#     elif page_selection == "🧾 OCA FINANCIAL":
#         financial_ui()
# # ... (the rest of your attendance tracker code continues here unchanged) ...



# ============================================================================= PART 2
# --------------------------------------------------------------------------
# Sidebar navigation
# --------------------------------------------------------------------------
# --------------------------------------------------------------------------
# Sidebar navigation
# --------------------------------------------------------------------------
def main():
    st.sidebar.title("📌 OCA Navigation")

    page_selection = st.sidebar.radio(
        "Select a Page:",
        ["🏠 Home", "🧾 Attendance Tracker", "🧾 OCA Finance","ACCOUNT API"],
        index=0
    )

    # ---------------------------------------------------------------
    # 📌 SIDEBAR: DELETE MEMBER FROM DATABASE (ALWAYS VISIBLE)
    # ---------------------------------------------------------------
    with st.sidebar.expander("🗑️ Delete Member From Database", expanded=False):

        st.markdown("Select a member to permanently delete:")

        members_df = st.session_state.get("members_df", pd.DataFrame())

        if not members_df.empty:

            member_to_delete = st.selectbox(
                "Member:",
                options=members_df["name"].tolist(),
                key="delete_member_sidebar"
            )

            if st.button("❌ Delete Selected Member", key="delete_member_button"):

                try:
                    import sqlite3
                    conn = sqlite3.connect("oca_attendance.db")
                    cur = conn.cursor()

                    # Delete attendance first
                    cur.execute("""
                        DELETE FROM attendance 
                        WHERE member_id = (SELECT id FROM members WHERE name = ?)
                    """, (member_to_delete,))

                    # Delete member
                    cur.execute("DELETE FROM members WHERE name = ?", (member_to_delete,))

                    conn.commit()
                    conn.close()

                    st.success(f"'{member_to_delete}' has been deleted successfully.")

                    # Refresh UI
                    updated = get_members()
                    st.session_state.members_df = updated
                    st.session_state.roster = updated["name"].sort_values().tolist()
                    _refresh_attendance_from_db()

                except Exception as e:
                    st.error(f"Error deleting member: {e}")

        else:
            st.info("No members in the database yet.")

    # ========================= DELETE ENDS HERE =========================

    # --------------------------------------------------------------------
    # HOME PAGE
    # --------------------------------------------------------------------
    # if page_selection == "🏠 Home":
    #     st.title("🏠 OCA Dashboard")
    #     st.markdown("""
    #     Welcome to the **Owerri Cultural Association Dashboard**
    #
    #     Navigate using the sidebar to access:
    #     - 🧾 Attendance Tracker
    #     - 🧾 OCA Financial System
    #     - 🔐 Account API
    #     """)

    # ====================================================================== FRONT PAGE LAYOUT

    # import streamlit as st

    # --------------------------------------------------------------------
    # HOME PAGE (OCA NorCal style landing)
    # --------------------------------------------------------------------
    if page_selection == "🏠 Home":
        st.markdown("""
        <style>
          .hero-wrap{
            border-radius: 18px;
            padding: 28px;
            border: 1px solid rgba(255,255,255,0.12);
            background: linear-gradient(135deg, rgba(16,185,129,0.18), rgba(59,130,246,0.18));
            margin-bottom: 16px;
          }
          .hero-kicker{
            display:inline-block;
            font-size: 0.85rem;
            padding: 6px 10px;
            border-radius: 999px;
            background: rgba(255,255,255,0.08);
            border: 1px solid rgba(255,255,255,0.12);
            margin-bottom: 10px;
          }
          .hero-title{
            margin: 0;
            line-height: 1.1;
            font-size: 2.1rem;
            font-weight: 800;
          }
          .hero-sub{
            margin-top: 10px;
            opacity: 0.9;
            font-size: 1.05rem;
          }
          .section-title{
            margin: 10px 0 8px 0;
            font-weight: 800;
            font-size: 1.1rem;
          }
          .card{
            border-radius: 16px;
            padding: 16px;
            border: 1px solid rgba(255,255,255,0.12);
            background: rgba(255,255,255,0.03);
            height: 100%;
          }
          .muted{opacity:0.85;}
          .divider{margin:16px 0; opacity:0.25;}
          .tag{
            display:inline-block;
            padding: 5px 10px;
            border-radius: 999px;
            background: rgba(255,255,255,0.06);
            border: 1px solid rgba(255,255,255,0.10);
            font-size: 0.85rem;
            margin-right: 8px;
            margin-top: 6px;
          }
        </style>
        """, unsafe_allow_html=True)

        # HERO
        st.markdown("""
        <div class="hero-wrap">
          <div class="hero-kicker">Owerri Cultural Association • Northern California</div>
          <div class="hero-title">Heritage. Unity. Strength.</div>
          <div class="hero-sub">
            OCA is a community of Owerri indigenes and families in Northern California—preserving culture,
            supporting one another, and building progress across generations.
          </div>
          <div style="margin-top:10px;">
            <span class="tag">Monthly Meetings</span>
            <span class="tag">Community Outreach</span>
            <span class="tag">Member Support</span>
            <span class="tag">Events & Celebrations</span>
          </div>
        </div>
        """, unsafe_allow_html=True)

        # PRIMARY CTA ROW
        c1, c2, c3 = st.columns([1.2, 1.2, 1])
        with c1:
            st.markdown('<div class="card"><div class="section-title">About OCA</div>'
                        '<div class="muted">We preserve Owerri traditions and promote unity, peace, and community development—'
                        'through meetings, events, and outreach.</div></div>', unsafe_allow_html=True)
        with c2:
            st.markdown('<div class="card"><div class="section-title">Our Automated System</div>'
                        '<div class="muted">This dashboard modernizes how we run OCA—attendance tracking, financial records, '
                        'member reporting, and secure admin tools—so leadership can operate faster, with clean records and transparency.</div></div>',
                        unsafe_allow_html=True)
        with c3:
            st.markdown('<div class="card"><div class="section-title">Quick Actions</div>'
                        '<div class="muted">Jump into the tools below.</div></div>', unsafe_allow_html=True)
            st.button("🧾 Attendance Tracker", use_container_width=True)
            st.button("💳 Financial System", use_container_width=True)
            st.button("🔐 Account / Admin API", use_container_width=True)

        st.markdown('<hr class="divider">', unsafe_allow_html=True)

        # PROGRAMS (like site sections)
        st.markdown('<div class="section-title">Our Programs</div>', unsafe_allow_html=True)
        p1, p2, p3 = st.columns(3)
        with p1:
            st.markdown("""
            <div class="card">
              <div class="section-title">🤝 Community Outreach</div>
              <div class="muted">Volunteer support, cultural impact, and assistance when members and families need it most.</div>
            </div>
            """, unsafe_allow_html=True)
        with p2:
            st.markdown("""
            <div class="card">
              <div class="section-title">🛡️ Member Support</div>
              <div class="muted">Structured support tied to participation and good standing—designed to keep our community strong.</div>
            </div>
            """, unsafe_allow_html=True)
        with p3:
            st.markdown("""
            <div class="card">
              <div class="section-title">🎉 Events & Culture</div>
              <div class="muted">Monthly meetings (Feb–Nov), annual celebrations, and gatherings that preserve identity and unity.</div>
            </div>
            """, unsafe_allow_html=True)

        st.markdown('<hr class="divider">', unsafe_allow_html=True)

        # UPCOMING / HIGHLIGHTS
        left, right = st.columns([1.3, 1])
        with left:
            st.markdown('<div class="section-title">Upcoming Events</div>', unsafe_allow_html=True)
            st.info(
                "Stay connected through meetings and community events. "
                "Use the Attendance module to confirm participation and view meeting history.",
                icon="📌"
            )
            with st.expander("Examples of what you can list here"):
                st.write("- Monthly Meeting (1st Sunday @ 6:00 PM)")
                st.write("- Annual OCA Events")
                st.write("- Community visits / outreach activities")

        with right:
            st.markdown('<div class="section-title">At a Glance</div>', unsafe_allow_html=True)
            m1, m2 = st.columns(2)
            with m1:
                st.metric("Years Active", "27+")
                st.metric("Members", "100+")
            with m2:
                st.metric("Villages Represented", "15+")
                st.metric("Doctorate Members", "5+")
            st.caption("Replace these with live metrics when your database is connected.")

        st.markdown('<hr class="divider">', unsafe_allow_html=True)

        # FOOTER
        st.markdown("""
        <div class="muted" style="padding: 6px 2px;">
          <b>OCA Dashboard</b> • Attendance • Finance • Secure Admin Tools<br/>
          Built to support transparency, continuity, and strong community operations.
        </div>
        """, unsafe_allow_html=True)











    # ///////////////////////////////////////////////////////////  ENDS HERE







    # --------------------------------------------------------------------
    # ATTENDANCE TRACKER PAGE
    # --------------------------------------------------------------------
    elif page_selection == "🧾 Attendance Tracker":
        attendance_ui()

    # --------------------------------------------------------------------
    # OCA FINANCIAL PAGE
    # --------------------------------------------------------------------
    elif page_selection == "🧾 OCA Finance":
        financial_ui()




    # --------------------------------------------------------------------
    # OCA FINANCIAL PAGE
    # --------------------------------------------------------------------
    elif page_selection == "ACCOUNT API":
        api_transactions_ui()




# ==================PART 2 ENDS HERE


# ================================================================================
# ACCOUNT API
# /////////////////////////////////////////////////////////////////////////////////

# ================================================================================
# ACCOUNT API
# /////////////////////////////////////////////////////////////////////////////////

# ================================================================================
# ACCOUNT API
# /////////////////////////////////////////////////////////////////////////////////

# # ======================================================================
# # BMO API — TRANSACTION VIEWER
# # ======================================================================
#
# import streamlit as st
# import requests
# # st.write("Loaded secrets:", st.secrets.keys())
# import pandas as pd
#
# def api_transactions_ui():
#
#     st.title("💳 BMO API — Transaction Viewer")
#     st.subheader("🔎 Search Account Transactions")
#     st.markdown("---")
#
#     # ------------------------------------------------------------------
#     # 1️⃣ GET BEARER TOKEN
#     # ------------------------------------------------------------------
#     def get_bearer_token():
#
#         url = "https://sandbox-open-api.bmo.com/oauth2/token"
#
#         data = {
#             "grant_type": "client_credentials",
#             "client_id": st.secrets["BMO_CLIENT_ID"],
#             "client_secret": st.secrets["BMO_CLIENT_SECRET"]
#         }
#
#         headers = {"Content-Type": "application/x-www-form-urlencoded"}
#
#         response = requests.post(url, data=data, headers=headers)
#
#         if response.status_code != 200:
#             st.error(f"❌ Token Error: {response.text}")
#             return None
#
#         return response.json().get("access_token")
#
#     # ------------------------------------------------------------------
#     # 2️⃣ FETCH TRANSACTIONS FROM BMO
#     # ------------------------------------------------------------------
#     def get_transactions(account_id, start_time, end_time, offset=0, limit=100):
#
#         token = get_bearer_token()
#         if not token:
#             return None
#
#         url = (
#             f"https://sandbox-open-api.bmo.com/open-banking/commercial-sb/accounts/"
#             f"{account_id}/transactions"
#             f"?startTime={start_time}&endTime={end_time}&offset={offset}&limit={limit}"
#         )
#
#         headers = {
#             "Accept": "application/json",
#             "Content-Type": "application/json",
#             "authorization": f"Bearer {token}",
#             "x-api-key": st.secrets["BMO_API_KEY"],
#             "distinct": "sybest_streamlit_v1"
#         }
#
#         response = requests.get(url, headers=headers)
#
#         if response.status_code != 200:
#             st.error(f"❌ API Error: {response.text}")
#             return None
#
#         return response.json()
#
#     # ------------------------------------------------------------------
#     # 3️⃣ USER INPUTS
#     # ------------------------------------------------------------------
#     account_id = st.text_input("Account ID", placeholder="Enter BMO Account ID")
#
#     col1, col2 = st.columns(2)
#     start_date = col1.date_input("Start Date")
#     end_date = col2.date_input("End Date")
#
#     start_ts = f"{start_date}T00:00:00Z"
#     end_ts = f"{end_date}T23:59:59Z"
#
#     # ------------------------------------------------------------------
#     # 4️⃣ FETCH BUTTON
#     # ------------------------------------------------------------------
#     if st.button("Fetch Transactions"):
#
#         if not account_id.strip():
#             st.warning("⚠️ Please enter an Account ID.")
#             return
#
#         data = get_transactions(account_id, start_ts, end_ts)
#
#         if not data:
#             return
#
#         tx_list = data.get("transactions", [])
#
#         if len(tx_list) == 0:
#             st.info("ℹ️ No transactions found for this date range.")
#             return
#
#         # ------------------------------------------------------------------
#         # 5️⃣ BUILD DATAFRAME
#         # ------------------------------------------------------------------
#         rows = []
#         for tx in tx_list:
#             d = tx["depositTransaction"]
#             rows.append({
#                 "Transaction ID": d["transactionId"],
#                 "Description": d["description"],
#                 "Amount": d["amount"],
#                 "Debit/Credit": d["debitCreditMemo"],
#                 "Posted Timestamp": d["postedTimestamp"],
#                 "Reference": d["bankReferenceNumber"]
#             })
#
#         df = pd.DataFrame(rows)
#
#         # ------------------------------------------------------------------
#         # 6️⃣ SUMMARY METRICS
#         # ------------------------------------------------------------------
#         st.markdown("### 📊 Summary")
#
#         colA, colB, colC = st.columns(3)
#
#         total_credits = df[df["Debit/Credit"] == "CREDIT"]["Amount"].sum()
#         total_debits = df[df["Debit/Credit"] == "DEBIT"]["Amount"].sum()
#         net_flow = df["Amount"].sum()
#
#         colA.metric("Total Credits", f"${total_credits:,.2f}")
#         colB.metric("Total Debits", f"${total_debits:,.2f}")
#         colC.metric("Net Cash Flow", f"${net_flow:,.2f}")
#
#         # ------------------------------------------------------------------
#         # 7️⃣ TABLE VIEW
#         # ------------------------------------------------------------------
#         st.markdown("### 📄 Transaction Details")
#         st.dataframe(df, use_container_width=True)
#
#         # ------------------------------------------------------------------
#         # 8️⃣ DOWNLOAD CSV
#         # ------------------------------------------------------------------
#         csv = df.to_csv(index=False).encode("utf-8")
#         st.download_button(
#             "⬇️ Download CSV",
#             csv,
#             "bmo_transactions.csv",
#             "text/csv",
#             use_container_width=True
#         )
#
#         st.success("✅ Transactions Loaded Successfully!")








# ======================================================================================================================================================
#                ATTENDANCE TRACKER CODE (your PART 444 section)
# ======================================================================================================================================================
# import streamlit as st
# import pandas as pd
# from datetime import date
# from io import BytesIO
# import matplotlib.pyplot as plt
#
# from db import add_member, get_members, upsert_attendance, get_attendance
#
#
# # # DATABASE ENDS HERE
#
#
# STATUS_OPTIONS = ["Present", "Absent", "Excused Absent"]
# ANNUAL_MEETINGS = 10
# MEETING_MONTHS = list(range(2, 12))
# MONTH_LABELS = ["FEB","MAR","APR","MAY","JUN","JUL","AUG","SEP","OCT","NOV"]
#
# # ---------------- Hardcoded roster ----------------
# DEFAULT_ROSTER = [
#     "Acholonu Kelechi & Nkiru (Mr. & Mrs.)",
#     "Akali Pet (Mrs.)",
#     "Amaechi Richard & Christina (Mr. & Mrs.)",
#     "Anyanwu Rosemary (Mrs)",
#     "Anyaso Hyacinth & Lydia (Chief & Lolo)",
#     "Chiagoro Chukwuma (Mr.)",
#     "Diala Emma & Joyce (Mr. & Mrs.)",
#     "Duru Gloria (Ms.)",
#     "Echibe Lucy (Ms.)",
#     "Egu Agnes (Mrs)",
#     "Egu Ernest (Mr. & Mrs.)",
#     "Egu Joy (Mrs)",
#     "Egu Mama Rose (Chief)",
#     "Egu Robert (Mr. & Mrs)",
#     "Eke Alexander (Mr. & Mrs.)",
#     "Eke Beatrice (Mrs.)",
#     "Eke Damian & Uloma (Chief & Lolo)",
#     "Ekeh George & Tina (Mr. & Mrs.)",
#     "Emeziem Adanma (Mrs.)",
#     "Ibe Prince & Chinwe (Mr. & Mrs.)",
#     "Ihejeto Anthony (Mr.)",
#     "Iwu Loretta (Ms.)",
#     "Iwuagwu Johnny (Mr. & Mrs.)",
#     "Mbara Tom (Mr. & Mrs.)",
#     "Njoku Stanley & Amaka (Mr. and Mrs)",
#     "Nwadibia Joannes Ubanwa (Ms.)",
#     "Nwadike Constance (Ms.)",
#     "Nwaogu Augustine (Mr. & Mrs.)",
#     "Nwaulu (Mr & Mrs)",
#     "Obichere Vivian (Mrs.)",
#     "Obilor Adolph (Chief & Dr. Lolo)",
#     "Obilor Innocent (Mr.)",
#     "Obilor Peter (Mr.)",
#     "Obinna Casmir (Mr.)",
#     "Ogbuehi Ethelbert (Mr.)",
#     "Ogwudire Kingsley & Ijeoma (Mr. & Mrs.)",
#     "Oha Linus & Liz (Mr. & Mrs.)",
#     "Oha Paul (Mr. & Mrs.)",
#     "Ohwobete Augustine and Grace (Dr. & Mrs.)",
#     "Ojibe Julie (Mrs)",
#     "Okoroafor Canon Progress (Rev)",
#     "Onyeagocha Rose (Lolo)",
#     "Onyejekwe Fidelia Chinyere (Mrs.)",
#     "Onyeneke Theophylus (Mr. & Mrs.)",
#     "Onyeukwu Damian (Chief)",
#     "Onyewuenyi Ethelbert R. (Mr.)",
#     "Opara Ugo Harris (Mr. & Mrs.)",
#     "Oparaocha Emeka & Joyce (Mr. & Mrs.)",
#     "Orji Christopher and Jacinta (Mr. & Mrs.)",
#     "Orji Ike F (Chief & Lolo)",
#     "Osuala Judy and Celine (Mr. & Mrs.)",
#     "Udeji Wilson",
#     "Ugorji Obi and Chinonye (Mr. & Mrs.)",
#     "Unaji Gibson & Victoria (Mr. & Mrs.)",
#     "Unanwa Christian (Dr. & Mrs)",
#     "Uwakwe Austin (Mr. & Mrs)",
#     "Uzohuo Uzoma David (Mr. & Mrs)",
#     "Uzoma Chinonye (Mr. & Mrs.)",
#     "Uzoma Christina (Ms.)",
#     "Yoko-Uzoma Okey (Mr. & Mrs.)",
# ]
#
# # ---------------- Session state ----------------
# def _init_state():
#     if "attendance_df" not in st.session_state:
#         st.session_state.attendance_df = pd.DataFrame(columns=["Date","Name","Status","Notes"])
#     if "roster" not in st.session_state:
#         st.session_state.roster = sorted(pd.unique(DEFAULT_ROSTER).tolist())
#
# # ---------------- Helpers ----------------
# def _upsert_record(att_date, name, status, note=""):
#     df = st.session_state.attendance_df
#     att_date = pd.to_datetime(att_date)
#     mask = (pd.to_datetime(df["Date"]) == att_date) & (df["Name"] == name)
#     if mask.any():
#         st.session_state.attendance_df.loc[mask, ["Status","Notes"]] = [status, note]
#     else:
#         st.session_state.attendance_df = pd.concat(
#             [df, pd.DataFrame([{"Date": att_date, "Name": name, "Status": status, "Notes": note}])],
#             ignore_index=True
#         )
#
# def _download_rollcall_csv(roster, att_date) -> bytes:
#     df = pd.DataFrame({
#         "Name": roster,
#         "Date": pd.to_datetime(att_date).date(),
#         "Attendance Status": ["" for _ in roster],
#         "Notes": ["" for _ in roster],
#     })
#     return df.to_csv(index=False).encode("utf-8")
#
# def _download_rollcall_xlsx(roster, att_date) -> bytes | None:
#     try:
#         out = BytesIO()
#         with pd.ExcelWriter(out, engine="openpyxl") as writer:
#             df = pd.DataFrame({
#                 "Name": roster,
#                 "Date": pd.to_datetime(att_date).date(),
#                 "Attendance Status": ["" for _ in roster],
#                 "Notes": ["" for _ in roster],
#             })
#             df.to_excel(writer, index=False, sheet_name="Roll Call")
#         return out.getvalue()
#     except Exception:
#         return None
#
# def _summary_cards(roster):
#     def _is_couple(n): return (" & " in n) or (" and " in n)
#     total_entries = len(roster)
#     couples = sum(1 for n in roster if _is_couple(n))
#     singles = total_entries - couples
#     total_individuals = couples*2 + singles
#     st.markdown(f"""
#         <style>
#             .cards {{display:flex; gap:16px; flex-wrap:wrap; margin: 4px 0 18px;}}
#             .card {{flex:1 1 240px; border:1px solid #e7e7e9; background:#fff; border-radius:16px;
#                     padding:16px 18px; box-shadow:0 2px 8px rgba(0,0,0,0.04);}}
#             .label {{font-size:13px; color:#666; margin-bottom:6px;}}
#             .value {{font-size:28px; font-weight:700; color:#111;}}
#         </style>
#         <div class="cards">
#           <div class="card"><div class="label">Total entries</div><div class="value">{total_entries}</div></div>
#           <div class="card"><div class="label">Singles</div><div class="value">{singles}</div></div>
#           <div class="card"><div class="label">Couples</div><div class="value">{couples}</div></div>
#           <div class="card"><div class="label">Total individual members</div><div class="value">{total_individuals}</div></div>
#         </div>
#     """, unsafe_allow_html=True)
#
# # ---------------- UI ----------------
# def attendance_ui():
#     _init_state()
#     st.title("🧾 Attendance — Easy Taker + Jan–Nov Report + Visuals")
#
#     _summary_cards(st.session_state.roster)
#
#     # 1) Roster
#     st.subheader("1) Roster & Add New Member")
#     c1, c2 = st.columns([2,1])
#     with c1:
#         new_name = st.text_input("Add a new member (full name)")
#         if st.button("➕ Add member"):
#             nm = new_name.strip()
#             if nm:
#                 if nm not in st.session_state.roster:
#                     st.session_state.roster.append(nm)
#                     st.session_state.roster = sorted(st.session_state.roster)
#                     st.success(f"Added: {nm}")
#                 else:
#                     st.info("That name already exists.")
#             else:
#                 st.warning("Enter a name.")
#     with c2:
#         roster_csv = pd.DataFrame({"Name": st.session_state.roster}).to_csv(index=False)
#         st.download_button("Download roster.csv", roster_csv, "roster.csv", "text/csv")
#
#     st.divider()
#     st.subheader("Comprehensive OCA Members List")
#     with st.expander("View Members List"):
#         st.dataframe(pd.DataFrame({"Name": st.session_state.roster}), use_container_width=True, hide_index=True)
#
#     st.divider()
#     st.subheader("2) Take Meeting Attendance")
#
#     att_date = st.date_input("Attendance date", value=date.today())
#     default_status = st.selectbox("Default for 'Mark All'", STATUS_OPTIONS, index=0)
#     if st.button("Mark All with default"):
#         for nm in st.session_state.roster:
#             _upsert_record(att_date, nm, default_status, "")
#         st.success("All marked with default value.")
#
#     st.divider()
#     st.download_button(
#         "Download roll-call CSV",
#         data=_download_rollcall_csv(st.session_state.roster, att_date),
#         file_name=f"rollcall_{att_date}.csv",
#         mime="text/csv"
#     )
#
#     xlsx_bytes = _download_rollcall_xlsx(st.session_state.roster, att_date)
#     if xlsx_bytes:
#         st.download_button(
#             "Download roll-call Excel",
#             data=xlsx_bytes,
#             file_name=f"rollcall_{att_date}.xlsx",
#             mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
#         )

# ====================================================================================================== PART 2
# # ===========================================================================================================
#
# ======================================================================================================================================================
#                ATTENDANCE TRACKER CODE (PART 444) — NOW USING SQLITE
# ======================================================================================================================================================

# ======================================================================================================================================================
#                ATTENDANCE TRACKER CODE (PART 444) — NOW USING SQLITE + UPLOAD + ANALYTICS
# ======================================================================================================================================================


import streamlit as st
import pandas as pd
from datetime import date, datetime
from io import BytesIO
import matplotlib.pyplot as plt



from db import add_member, get_members, upsert_attendance, get_attendance, update_member_name


STATUS_OPTIONS = ["Present", "Absent", "Excused Absent"]
ANNUAL_MEETINGS = 10
MEETING_MONTHS = list(range(2, 12))
MONTH_LABELS = ["FEB","MAR","APR","MAY","JUN","JUL","AUG","SEP","OCT","NOV"]

# ---------------- Initial hardcoded roster (used to seed DB on first run) ----------------
DEFAULT_ROSTER = [
    "Acholonu Kelechi & Nkiru (Mr. & Mrs.)",
    "Akali Pet (Mrs.)",
    "Amaechi Richard & Christina (Mr. & Mrs.)",
    "Anyanwu Rosemary (Mrs)",
    "Anyaso Hyacinth & Lydia (Chief & Lolo)",
    "Chiagoro Chukwuma (Mr.)",
    "Diala Emma & Joyce (Mr. & Mrs.)",
    "Duru Gloria (Ms.)",
    "Echibe Lucy (Ms.)",
    "Egu Agnes (Mrs)",
    "Egu Ernest (Mr. & Mrs.)",
    "Egu Joy (Mrs)",
    "Egu Mama Rose (Chief)",
    "Egu Robert (Mr. & Mrs)",
    "Eke Alexander (Mr. & Mrs.)",
    "Eke Beatrice (Mrs.)",
    "Eke Damian & Uloma (Chief & Lolo)",
    "Ekeh George & Tina (Mr. & Mrs.)",
    "Emeziem Adanma (Mrs.)",
    "Ibe Prince & Chinwe (Mr. & Mrs.)",
    "Ihejeto Anthony (Mr.)",
    "Iwu Loretta (Ms.)",
    "Iwuagwu Johnny (Mr. & Mrs.)",
    "Mbara Thomas (Mr. & Mrs.)",
    "Njoku Stanley & Amaka (Dr. and Mrs)",
    "Nwadibia Joannes Ubanwa (Ms.)",
    "Nwadike Constance (Ms.)",
    "Nwaogu Augustine (Mr. & Mrs.)",
    "Nwaulu (Mr & Mrs)",
    "Obichere Vivian (Mrs.)",
    "Obilor Adolph (Chief & Dr. Lolo)",
    "Obilor Innocent (Mr.)",
    "Obilor Peter (Mr.)",
    "Obinna Casmir (Mr.)",
    "Ogbuehi Ethelbert (Mr.)",
    "Ogwudire Kingsley & Ijeoma (Mr. & Mrs.)",
    "Oha Linus & Liz (Mr. & Mrs.)",
    "Oha Paul (Mr. & Mrs.)",
    "Ohwobete Augustine and Grace (Dr. & Mrs.)",
    "Ojibe Julie (Mrs)",
    "Okoroafor Canon Progress (Rev)",
    "Onyeagocha Rose (Lolo)",
    "Onyejekwe Fidelia Chinyere (Mrs.)",
    "Onyeneke Theophylus (Mr. & Mrs.)",
    "Onyeukwu Damian (Chief)",
    "Onyewuenyi Ethelbert R. (Mr.)",
    "Opara Ugo Harris (Mr. & Mrs.)",
    "Oparaocha Emeka & Joyce (Mr. & Mrs.)",
    "Orji Christopher and Jacinta (Mr. & Mrs.)",
    "Orji Ike F (Chief & Lolo)",
    "Osuala Judy and Celine (Mr. & Mrs.)",
    "Udeji Wilson",
    "Ugorji Obi and Chinonye (Mr. & Mrs.)",
    "Unaji Gibson & Victoria (Mr. & Mrs.)",
    "Uwakwe Austin (Mr. & Mrs)",
    "Uzohuo Uzoma David (Mr. & Mrs)",
    "Uzoma Chinonye (Mr. & Mrs.)",
    "Uzoma Christina (Ms.)",
    "Yoko-Uzoma Okey (Mr. & Mrs.)",
]

# ---------------- Session state (now backed by DB) ----------------
def _init_state():
    # Members / roster from DB (seed from DEFAULT_ROSTER first time)
    if "members_df" not in st.session_state or "roster" not in st.session_state:
        members_df = get_members()
        if members_df.empty:
            # Seed database the very first time
            for nm in DEFAULT_ROSTER:
                add_member(nm)
            members_df = get_members()
        st.session_state.members_df = members_df
        st.session_state.roster = members_df["name"].sort_values().tolist()

        st.session_state.numbered_roster = [
            f"({i + 1}) {name}" for i, name in enumerate(st.session_state.roster)   #TO ADD NUMBERING AFTYER CHURCH
        ]


    # Attendance history from DB (for visualization / table)
    if "attendance_df" not in st.session_state:
        att_df = get_attendance()
        if not att_df.empty:
            att_df = att_df.rename(
                columns={"date": "Date", "name": "Name", "status": "Status", "notes": "Notes"}
            )
        else:
            att_df = pd.DataFrame(columns=["Date", "Name", "Status", "Notes"])
        st.session_state.attendance_df = att_df

# ---------------- Helpers ----------------
def _refresh_attendance_from_db():
    att_df = get_attendance()
    if not att_df.empty:
        att_df = att_df.rename(
            columns={"date": "Date", "name": "Name", "status": "Status", "notes": "Notes"}
        )
        att_df["Date"] = pd.to_datetime(att_df["Date"])
    else:
        att_df = pd.DataFrame(columns=["Date", "Name", "Status", "Notes"])
    st.session_state.attendance_df = att_df

def _upsert_record(att_date, name, status, note=""):
    """Save/update one record in SQLite AND keep session dataframe in sync."""
    if pd.isna(att_date) or not str(name).strip() or str(status) not in STATUS_OPTIONS:
        return

    # Ensure date is ISO string for DB
    att_date_iso = pd.to_datetime(att_date).date().isoformat()

    # Find member_id from name (or add if somehow missing)
    members_df = st.session_state.members_df
    row = members_df.loc[members_df["name"] == name]
    if row.empty:
        add_member(name)
        members_df = get_members()
        st.session_state.members_df = members_df
        st.session_state.roster = members_df["name"].sort_values().tolist()
        row = members_df.loc[members_df["name"] == name]
    member_id = int(row.iloc[0]["id"])

    # Write to DB
    upsert_attendance(att_date_iso, member_id, status, note)

    # Update in-memory df for immediate UI use
    df = st.session_state.attendance_df
    d = pd.to_datetime(att_date_iso)
    mask = (pd.to_datetime(df["Date"]) == d) & (df["Name"] == name)
    if mask.any():
        st.session_state.attendance_df.loc[mask, ["Status", "Notes"]] = [status, note]
    else:
        st.session_state.attendance_df = pd.concat(
            [df, pd.DataFrame([{"Date": d, "Name": name, "Status": status, "Notes": note}])],
            ignore_index=True,
        )

def _download_rollcall_csv(roster, att_date) -> bytes:
    df = pd.DataFrame({
        "Name": roster,
        "Date": pd.to_datetime(att_date).date(),
        "Attendance Status": ["" for _ in roster],
        "Notes": ["" for _ in roster],
    })
    return df.to_csv(index=False).encode("utf-8")

def _download_rollcall_xlsx(roster, att_date) -> bytes | None:
    try:
        out = BytesIO()
        with pd.ExcelWriter(out, engine="openpyxl") as writer:
            df = pd.DataFrame({
                "Name": roster,
                "Date": pd.to_datetime(att_date).date(),
                "Attendance Status": ["" for _ in roster],
                "Notes": ["" for _ in roster],
            })
            df.to_excel(writer, index=False, sheet_name="Roll Call")
        return out.getvalue()
    except Exception:
        return None

def _download_attendance_template() -> bytes:
    """
    Template for importing attendance via CSV/Excel.
    Required columns: Date, Name, Status, Notes
    """
    df = pd.DataFrame({
        "Date": [date.today().isoformat()],
        "Name": ["Example Member"],
        "Status": ["Present"],  # Must be one of STATUS_OPTIONS
        "Notes": [""],
    })
    return df.to_csv(index=False).encode("utf-8")

def _summary_cards(roster):
    def _is_couple(n): return (" & " in n) or (" and " in n)
    total_entries = len(roster)
    couples = sum(1 for n in roster if _is_couple(n))
    singles = total_entries - couples
    total_individuals = couples * 2 + singles
    st.markdown(f"""
        <style>
            .cards {{display:flex; gap:16px; flex-wrap:wrap; margin: 4px 0 18px;}}
            .card {{flex:1 1 240px; border:1px solid #e7e7e9; background:#fff; border-radius:16px;
                    padding:16px 18px; box-shadow:0 2px 8px rgba(0,0,0,0.04);}}
            .label {{font-size:13px; color:#666; margin-bottom:6px;}}
            .value {{font-size:28px; font-weight:700; color:#111;}}
        </style>
        <div class="cards">
          <div class="card"><div class="label">Total entries</div><div class="value">{total_entries}</div></div>
          <div class="card"><div class="label">Singles</div><div class="value">{singles}</div></div>
          <div class="card"><div class="label">Couples</div><div class="value">{couples}</div></div>
          <div class="card"><div class="label">Total individual members</div><div class="value">{total_individuals}</div></div>
        </div>
    """, unsafe_allow_html=True)

# ========================================================== COLOR CODE UPDATED
# ===================================================================
# COLOR-CODED DROPDOWN FOR ATTENDANCE SELECTBOX (CSS HACK)
# ===================================================================

def colored_selectbox(label, key, default=None):
    """
    Custom colored selectbox using HTML + CSS trick.
    This replaces st.selectbox visually.
    """
    color_map = {
        "Present": "#2ecc71",        # green
        "Absent": "#e74c3c",         # red
        "Excused Absent": "#3498db"  # blue
    }

    # Create a colored dropdown using HTML <select>
    html = f"""
    <style>
        select {{
            padding: 8px 12px;
            border-radius: 6px;
            font-size: 15px;
            background-color: #111;
            color: white;
            border: 1px solid #333;
        }}
        option[value="Present"] {{
            color: {color_map["Present"]};
            font-weight: 600;
        }}
        option[value="Absent"] {{
            color: {color_map["Absent"]};
            font-weight: 600;
        }}
        option[value="Excused Absent"] {{
            color: {color_map["Excused Absent"]};
            font-weight: 600;
        }}
    </style>

    <label style="font-size:14px; color:#ccc;">{label}</label>
    <select id="{key}">
        <option value="Present" {"selected" if default=="Present" else ""}>Present</option>
        <option value="Absent" {"selected" if default=="Absent" else ""}>Absent</option>
        <option value="Excused Absent" {"selected" if default=="Excused Absent" else ""}>Excused Absent</option>
    </select>

    <script>
        const selectElem = document.getElementById("{key}");
        selectElem.addEventListener("change", function() {{
            window.parent.postMessage({{ key: "{key}", value: selectElem.value }}, "*");
        }});
    </script>
    """

    st.markdown(html, unsafe_allow_html=True)

    # Get the value from session_state
    return st.session_state.get(key, default)

    # return st.session_state.get(key, "Absent")  # TO MAKE ABSENT DEFAULT UPDATED

# ========================================================================= CSS COLOR CODE
st.markdown("""
<style>

div[data-baseweb="select"] span {
    font-weight: 600 !important;
}

/* Selected item background + text */
div[data-baseweb="select"] div[role="button"] span {
    padding: 4px 8px !important;
    border-radius: 6px !important;
}

/* ---------- PRESENT (GREEN SHIELD) ---------- */
div[data-baseweb="select"] div[role="listbox"] div[data-option*="Present"] {
    background-color: rgba(46, 204, 113, 0.25) !important;
    color: #2ecc71 !important;
    font-weight: 700 !important;
    border-radius: 6px;
}
div[data-baseweb="select"] div[role="button"] span:has-text("Present") {
    background-color: rgba(46, 204, 113, 0.30) !important;
    color: #2ecc71 !important;
}

/* ---------- ABSENT (RED SHIELD) ---------- */
div[data-baseweb="select"] div[role="listbox"] div[data-option*="Absent"] {
    background-color: rgba(231, 76, 60, 0.25) !important;
    color: #e74c3c !important;
    font-weight: 700 !important;
    border-radius: 6px;
}
div[data-baseweb="select"] div[role="button"] span:has-text("Absent") {
    background-color: rgba(231, 76, 60, 0.30) !important;
    color: #e74c3c !important;
}

/* ----- EXCUSED ABSENT (BLUE SHIELD) ----- */
div[data-baseweb="select"] div[role="listbox"] div[data-option*="Excused"] {
    background-color: rgba(52, 152, 219, 0.25) !important;
    color: #3498db !important;
    font-weight: 700 !important;
    border-radius: 6px;
}
div[data-baseweb="select"] div[role="button"] span:has-text("Excused") {
    background-color: rgba(52, 152, 219, 0.30) !important;
    color: #3498db !important;
}

</style>
""", unsafe_allow_html=True)








# =========================== COLOR CODE ENDED HERE


# ---------------- UI ----------------
# def attendance_ui():
#     _init_state()
#     st.title("🧾 Attendance — Easy Taker + Upload + Analytics (DB Saved)")

# ======================================================================================= PART2
def attendance_ui():
    _init_state()

    col1, col2 = st.columns([1, 5])

    with col1:
        logo_path = r"C:/Users/stans/OneDrive/Desktop/OCA/OCA LOGO/OCA LOGO.JPG"
        st.image(logo_path, width=120)

    with col2:
        st.title("🧾 OCA Attendance & Analytics Dashboard")
    st.divider()
# ====================================================================================== LOGO ADDED


    _summary_cards(st.session_state.roster)

    # 1) Roster
    st.subheader("1) Roster: Add / Edit Members")
    c1, c2, c3 = st.columns([2, 2, 1])

    # Add member
    with c1:
        new_name = st.text_input("Add a new member (Formate Last, First (Title )")
        if st.button("➕ Add member"):
            nm = new_name.strip()
            if nm:
                if nm not in st.session_state.roster:
                    add_member(nm)
                    members_df = get_members()
                    st.session_state.members_df = members_df
                    st.session_state.roster = members_df["name"].sort_values().tolist()
                    st.success(f"Added: {nm}")
                else:
                    st.info("That name already exists.")
            else:
                st.warning("Enter a name.")

    # Edit member name
    with c2:
        members_df = st.session_state.members_df
        if not members_df.empty:
            edit_target = st.selectbox(
                "Select member to rename",
                options=members_df["name"].tolist()
            )
            new_label = st.text_input("New name for selected member")
            if st.button("✏️ Rename member"):
                if new_label.strip():
                    row = members_df.loc[members_df["name"] == edit_target]
                    if not row.empty:
                        member_id = int(row.iloc[0]["id"])
                        ok = update_member_name(member_id, new_label.strip())
                        if ok:
                            st.success(f"Renamed '{edit_target}' → '{new_label.strip()}'")
                            members_df = get_members()
                            st.session_state.members_df = members_df
                            st.session_state.roster = members_df["name"].sort_values().tolist()
                            # Also refresh attendance names
                            _refresh_attendance_from_db()
                        else:
                            st.error("A member with that name already exists.")
                else:
                    st.warning("Enter a new name.")
        else:
            st.info("No members in the database yet.")

    # Download roster
    with c3:
        roster_csv = pd.DataFrame({"Name": st.session_state.roster}).to_csv(index=False)
        st.download_button("Download roster.csv", roster_csv, "roster.csv", "text/csv")

    st.divider()
    st.subheader("Comprehensive OCA Members List")
    with st.expander("View Members List"):
        st.dataframe(pd.DataFrame({"Name": st.session_state.roster}),
                     use_container_width=True, hide_index=True)


# ======================================================================================== STEP 1 ENDS HERE
# ======================================================================================== STEP 1 ENDS HERE
# ======================================================================================== STEP 1 ENDS HERE


    # 2) Take Meeting Attendance
    st.divider()
    st.subheader("2) Take Meeting Attendance")

    # ============================
    # MODE SELECTOR
    # ============================
    mode = st.radio(
        "Select Mode:",
        ["Take New Attendance", "Edit Existing Attendance"],
        horizontal=True,
        index=0
    )

    # ============================
    # MODE 1: TAKE NEW ATTENDANCE
    # ============================
    if mode == "Take New Attendance":

        att_date = st.date_input("Attendance date", value=date.today())

        default_status = st.selectbox(
            "Default for 'Mark All'",
            STATUS_OPTIONS,
            index=0
        )

        if st.button("Mark All with default"):
            for nm in st.session_state.roster:
                _upsert_record(att_date, nm, default_status, "")
            st.success("All marked with default value and saved to database.")
            _refresh_attendance_from_db()

        st.divider()
        st.markdown("### Update individual attendance & notes (NEW)")

    # ============================
    # MODE 2: EDIT EXISTING ATTENDANCE
    # ============================
    else:

        st.markdown("### 🔄 Edit Existing Attendance Records")

        # Load all saved dates from DB
        all_dates = sorted(
            st.session_state.attendance_df["Date"].dt.date.unique().tolist()
        ) if not st.session_state.attendance_df.empty else []

        if all_dates:
            edit_date = st.selectbox(
                "Select date to edit attendance",
                options=all_dates,
                key="edit_attendance_date"
            )

            st.info(f"Loaded attendance for: **{edit_date}**")

            # Filter attendance for selected date
            df_edit = st.session_state.attendance_df[
                st.session_state.attendance_df["Date"].dt.date == edit_date
                ]

            # Lookup dictionaries
            saved_status = {row["Name"]: row["Status"] for _, row in df_edit.iterrows()}
            saved_notes = {row["Name"]: row["Notes"] for _, row in df_edit.iterrows()}

        else:
            st.warning("No previous attendance records found.")
            edit_date = None
            saved_status = {}
            saved_notes = {}

        st.divider()
        st.markdown("### Update individual attendance & notes (EDIT)")

    # ---------------------------------------------------------
    # NOW BEGIN THE LOOP THAT HANDLES BOTH MODES
    # ---------------------------------------------------------

    # for nm in st.session_state.roster:

    for idx, display_name in enumerate(st.session_state.numbered_roster):  #ADDED NAME PARENTENISS
        nm = st.session_state.roster[idx]  # REAL name (saved to DB)

        # Load values depending on mode
        if mode == "Edit Existing Attendance":
            existing_status = saved_status.get(nm, "Absent")
            existing_note = saved_notes.get(nm, "")
            date_to_save = edit_date
        else:
            existing_status = default_status
            existing_note = ""
            date_to_save = att_date

        c1, c2, c3 = st.columns([2.5, 1.2, 3])

        # ===============================================================All ENDS HERE

        # --- NAME ---
        with c1:
            st.markdown(f"<div style='padding-top:8px;'>{nm}</div>", unsafe_allow_html=True)

        # --- STATUS DROPDOWN INLINE ---
        # --- STATUS DROPDOWN INLINE ---
        with c2:
            col_status, col_badge = st.columns([1.4, 0.8])

            # LEFT: Dropdown
            with col_status:
                status = st.selectbox(
                    "",
                    STATUS_OPTIONS,
                    index=STATUS_OPTIONS.index(existing_status),
                    key=f"status_{nm}_{date_to_save}",
                    label_visibility="collapsed"
                )

            # RIGHT: Color badge
            with col_badge:
                color_map = {
                    "Present": "#27ae60",
                    "Absent": "#e74c3c",
                    "Excused Absent": "#2980b9",
                }

                badge_html = f"""
                <div style="
                    margin-top:4px;
                    padding:4px 10px;
                    border-radius:20px;
                    background:{color_map[status]};
                    color:white;
                    font-size:11px;
                    font-weight:600;
                    text-align:center;
                    white-space:nowrap;">
                    {status}
                </div>
                """

                st.markdown(badge_html, unsafe_allow_html=True)

        # --- NOTES INLINE ---
        with c3:
            note = st.text_input(
                "",
                value=existing_note,
                key=f"note_{nm}_{date_to_save}",
                placeholder=f"Note for {nm}",
                label_visibility="collapsed"
            )

        # Save record
        _upsert_record(date_to_save, nm, status, note)


    # ========================================================================== END FOOD
    # ========================================================================== END FOOD

    st.success("Attendance saved/updated in database for selected date.")
    _refresh_attendance_from_db()

    st.success("Attendance saved/updated in database.")

    # ============================================================
    # SAFE DATE FOR DOWNLOAD BUTTONS (FIXES att_date ERROR)
    # ============================================================

    try:
        selected_att_date = att_date
    except NameError:
        try:
            selected_att_date = edit_date
        except NameError:
            selected_att_date = date.today()

    # ============================================================
    # 3) Download Roll-call & Import Template
    # ============================================================

    st.divider()
    st.subheader("3) Download Roll-call & Import Template")

    col_csv, col_xlsx, col_tpl = st.columns(3)
    with col_csv:
        st.download_button(
            "Download roll-call CSV",
            data=_download_rollcall_csv(st.session_state.roster, selected_att_date),
            file_name=f"rollcall_{selected_att_date}.csv",
            mime="text/csv"
        )

    with col_xlsx:
        xlsx_bytes = _download_rollcall_xlsx(st.session_state.roster, selected_att_date)
        if xlsx_bytes:
            st.download_button(
                "Download roll-call Excel",
                data=xlsx_bytes,
                file_name=f"rollcall_{selected_att_date}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

    with col_tpl:
        st.download_button(
            "Attendance Import Template (CSV)",
            data=_download_attendance_template(),
            file_name="attendance_template.csv",
            mime="text/csv",
            help="Columns: Date, Name, Status, Notes. Status must be Present / Absent / Excused Absent."
        )


    # 4) Import attendance from CSV/Excel
    st.divider()
    st.subheader("4) Upload Attendance (CSV / Excel → Save to DB)")

    uploaded = st.file_uploader("Upload attendance file (CSV or Excel)", type=["csv", "xlsx"])
    if uploaded is not None:
        try:
            if uploaded.name.lower().endswith(".csv"):
                df_up = pd.read_csv(uploaded)
            else:
                df_up = pd.read_excel(uploaded)

            required_cols = {"Date", "Name", "Status"}
            if not required_cols.issubset(set(df_up.columns)):
                st.error("File must contain at least these columns: Date, Name, Status. Optional: Notes.")
            else:
                df_up["Notes"] = df_up.get("Notes", "")

                bad_status = sorted(set(df_up["Status"].dropna()) - set(STATUS_OPTIONS))
                if bad_status:
                    st.warning(f"Ignoring rows with invalid Status values: {bad_status}")

                imported_rows = 0
                for _, r in df_up.iterrows():
                    if pd.isna(r["Date"]) or pd.isna(r["Name"]) or pd.isna(r["Status"]):
                        continue
                    if r["Status"] not in STATUS_OPTIONS:
                        continue
                    _upsert_record(r["Date"], str(r["Name"]), str(r["Status"]), str(r.get("Notes", "") or ""))
                    imported_rows += 1

                _refresh_attendance_from_db()
                st.success(f"Imported / updated {imported_rows} attendance rows into database.")
        except Exception as e:
            st.error(f"Error reading uploaded file: {e}")

    # 5) Attendance history + visuals + analytics
    st.divider()
    st.subheader("5) Attendance History, Graphs & Missed-Meetings Tracker")

    att_df = st.session_state.attendance_df
    if not att_df.empty:
        st.markdown("##### Full Attendance Table")
        st.dataframe(att_df.sort_values("Date", ascending=False),
                     use_container_width=True, hide_index=True)

        # # --- Year filter ---
        # current_year = datetime.now().year
        # years_available = sorted(att_df["Date"].dt.year.unique())
        # default_year = current_year if current_year in years_available else years_available[-1]
        # year_selected = st.selectbox("Select year for analysis", options=years_available, index=years_available.index(default_year))

        # =========================== PART 2

        # ========================================================================
        # YEAR FILTER (FULL FIXED VERSION)
        # ========================================================================

        from datetime import datetime

        # Make sure Date column is in datetime format
        att_df["Date"] = pd.to_datetime(att_df["Date"], errors="coerce")

        # Remove rows with invalid dates
        att_df = att_df.dropna(subset=["Date"])

        # If nothing left after cleaning → stop
        if att_df.empty:
            st.info("No valid attendance records found.")
            st.stop()

        # Extract available years
        att_df["Year"] = att_df["Date"].dt.year
        years_available = sorted(att_df["Year"].unique().tolist())

        if not years_available:
            st.info("No years available for filtering.")
            st.stop()

        # Determine default year
        current_year = datetime.now().year
        default_year = current_year if current_year in years_available else years_available[-1]

        # UI Selectbox for year selection
        year_selected = st.selectbox(
            "📅 Select Year for Attendance Analysis",
            options=years_available,
            index=years_available.index(default_year)
        )

        # Filter to selected year
        df_year = att_df[att_df["Year"] == year_selected].copy()

        # If the selected year has no attendance
        if df_year.empty:
            st.warning(f"No attendance records found for {year_selected}.")
            st.stop()

        # Convert Month number
        df_year["Month"] = df_year["Date"].dt.month

        # Provide cleaned, year-filtered dataframe for further charts/tables
        st.success(f"Showing attendance for YEAR {year_selected}")

        # ============================ PART 2 ENDS HERE  FILTER MONTHS STARTS HERE

        # ========================================================================
        # OPTIONAL — MONTH FILTER (BAR, PIE, SUMMARY READY)
        # ========================================================================

        month_labels_map = {
            1: "JAN", 2: "FEB", 3: "MAR", 4: "APR", 5: "MAY", 6: "JUN",
            7: "JUL", 8: "AUG", 9: "SEP", 10: "OCT", 11: "NOV", 12: "DEC"
        }

        months_available = sorted(df_year["Month"].unique().tolist())
        month_options = [month_labels_map[m] for m in months_available]

        month_selected_label = st.selectbox(
            "📅 Select Month (Optional)",
            ["ALL"] + month_options
        )

        if month_selected_label != "ALL":
            month_selected = {v: k for k, v in month_labels_map.items()}[month_selected_label]
            df_filtered = df_year[df_year["Month"] == month_selected].copy()
            st.success(f"Showing data for {month_selected_label} {year_selected}")
        else:
            df_filtered = df_year.copy()
            st.success(f"Showing data for ALL months in {year_selected}")

        # ============================================== ENDS HERE




        df_year = att_df[att_df["Date"].dt.year == year_selected].copy()

        if not df_year.empty:
            # Status totals
            status_counts = df_year["Status"].value_counts().reindex(STATUS_OPTIONS, fill_value=0)
            c1, c2, c3 = st.columns(3)
            c1.metric("Total Present", int(status_counts["Present"]))
            c2.metric("Total Absent", int(status_counts["Absent"]))
            c3.metric("Total Excused Absent", int(status_counts["Excused Absent"]))

            with st.expander("See who was Present / Absent / Excused in selected year"):
                col_a, col_b, col_c = st.columns(3)
                with col_a:
                    st.markdown("**Present (unique members)**")
                    st.write(sorted(df_year[df_year["Status"]=="Present"]["Name"].unique().tolist()))
                with col_b:
                    st.markdown("**Absent (unique members)**")
                    st.write(sorted(df_year[df_year["Status"]=="Absent"]["Name"].unique().tolist()))
                with col_c:
                    st.markdown("**Excused Absent (unique members)**")
                    st.write(sorted(df_year[df_year["Status"]=="Excused Absent"]["Name"].unique().tolist()))

# ====================================================================================================================
# GRAPHS INSIGHTS
# ==================================================================================================================

            # Graph: stacked bar by Name & Status
            st.markdown("##### Per-member Attendance Breakdown (stacked)")

            pivot = pd.crosstab(
                df_year["Name"], df_year["Status"]
            ).reindex(columns=STATUS_OPTIONS, fill_value=0)

            if not pivot.empty:
                fig1, ax1 = plt.subplots(figsize=(12, 5))

                pivot.plot(
                    kind="bar",
                    stacked=True,
                    ax=ax1,
                    color=["#27ae60", "#e74c3c", "#2980b9"]  # Present, Absent, Excused
                )

                ax1.set_ylabel("Meetings")
                ax1.set_title(f"Attendance Breakdown by Member — {year_selected}")
                plt.xticks(rotation=90)

                # ---- FIX LEGEND POSITION ----
                ax1.legend(
                    title="Status",
                    loc="upper left",
                    bbox_to_anchor=(1.02, 1),
                    frameon=True
                )

                plt.tight_layout()
                st.pyplot(fig1)



# =========================================================================================== PART 2


            # ====================================================================== ADDON MONTHLY BLOCK


            # =====================================================================
            # GROUPED MONTHLY ATTENDANCE (Present, Absent, Excused Absent)
            # =====================================================================

            df_year["MonthNum"] = df_year["Date"].dt.month

            # Count each attendance status per month
            monthly_counts = df_year.groupby(["MonthNum", "Status"])["Name"].nunique().unstack(fill_value=0)

            # Ensure all 3 columns exist
            for col in STATUS_OPTIONS:
                if col not in monthly_counts.columns:
                    monthly_counts[col] = 0

            # Force correct column order for consistent graph colors
            monthly_counts = monthly_counts[["Present", "Absent", "Excused Absent"]]

            # Convert month numbers to consistent labels (JAN–DEC)
            monthly_counts.index = [month_labels_map.get(m, str(m)) for m in monthly_counts.index]

            st.markdown("##### 📊 Monthly Attendance — Present, Absent & Excused Absent")

            fig3, ax3 = plt.subplots(figsize=(10, 5))

            # COLOR MAP (MATCHING YOUR APP)
            color_map = {
                "Present": "#27ae60",  # green
                "Absent": "#e74c3c",  # red
                "Excused Absent": "#2980b9"  # blue
            }

            # Plot grouped bars with correct color alignment
            monthly_counts.plot(
                kind="bar",
                ax=ax3,
                color=[color_map["Present"], color_map["Absent"], color_map["Excused Absent"]],
                width=0.75
            )

            # ANNOTATE ALL BARS
            for container in ax3.containers:
                for bar in container:
                    height = bar.get_height()
                    if height > 0:
                        ax3.annotate(
                            str(int(height)),
                            xy=(bar.get_x() + bar.get_width() / 2, height),
                            xytext=(0, 3),
                            textcoords="offset points",
                            ha="center",
                            fontsize=9,
                            color="#111"
                        )

            ax3.set_ylabel("Number of Members")
            ax3.set_title(f"Monthly Attendance Breakdown — {year_selected}", fontsize=13)
            ax3.legend(title="Status", fontsize=9)
            plt.xticks(rotation=0)
            plt.tight_layout()

            st.pyplot(fig3)

            # ================================================== ENDS HERE ADDON MONTHLY BLOCK


            # =====================================================================
            # 🚨 TRACKER: Combined Absence Report (Absent + Excused Absent)
            # =====================================================================

            st.markdown("### 🚨 Attendance Concern Report — Members With 5+ Total Missed Meetings")

            # Month labels lookup
            month_labels_map = {
                1: "JAN", 2: "FEB", 3: "MAR", 4: "APR", 5: "MAY", 6: "JUN",
                7: "JUL", 8: "AUG", 9: "SEP", 10: "OCT", 11: "NOV", 12: "DEC"
            }

            # Add MonthNum and MonthLabel
            df_year["MonthNum"] = df_year["Date"].dt.month
            df_year["MonthLabel"] = df_year["MonthNum"].map(month_labels_map)

            # Filter only missed statuses
            missed = df_year[df_year["Status"].isin(["Absent", "Excused Absent"])].copy()

            if missed.empty:
                st.info("No absence records for the selected year.")
            else:
                # Summary table
                summary = (
                    missed.groupby("Name")
                    .agg(
                        TotalMissed=("Status", "size"),
                        AbsentCount=("Status", lambda x: (x == "Absent").sum()),
                        ExcusedCount=("Status", lambda x: (x == "Excused Absent").sum()),
                        MissedMonths=("MonthLabel", lambda x: ", ".join(sorted(set(x))))
                    )
                    .reset_index()
                )

                # Filter members with 5+ total missed
                flagged = summary[summary["TotalMissed"] >= 5] \
                    .sort_values("TotalMissed", ascending=False)

                # -------------------------------------------
                # ⭐ SHOW TOTAL NUMBER OF MEMBERS WHO MISSED 5+
                # -------------------------------------------
                total_flagged = len(flagged)
                st.metric("Total Members With 5+ Missed Meetings", total_flagged)

                if flagged.empty:
                    st.info("No member has 5 or more total missed meetings.")
                else:
                    st.dataframe(
                        flagged,
                        use_container_width=True,
                        hide_index=True
                    )

                    # Detailed breakdown
                    with st.expander("📌 Detailed Insights for Members With 5+ Absences"):
                        for _, row in flagged.iterrows():
                            st.markdown(
                                f"""
                                **{row['Name']}**  
                                - **Total Missed:** {row['TotalMissed']}  
                                - **Absent:** {row['AbsentCount']}  
                                - **Excused Absent:** {row['ExcusedCount']}  
                                - **Months missed:** {row['MissedMonths']}  
                                """
                            )


# =================================================================== ATTENDANCE ENDS HERE

                # =====================================================================
                # 📊 ADDITIONAL INSIGHTS: Top Excused & Unexcused Absences
                # =====================================================================

                st.markdown("---")
                st.markdown("## 📊 Additional Attendance Insights")

                # Ensure df_year already has MonthLabel
                if "MonthLabel" not in df_year.columns:
                    df_year["MonthNum"] = df_year["Date"].dt.month
                    month_labels_map = {
                        1: "JAN", 2: "FEB", 3: "MAR", 4: "APR", 5: "MAY", 6: "JUN",
                        7: "JUL", 8: "AUG", 9: "SEP", 10: "OCT", 11: "NOV", 12: "DEC"
                    }
                    df_year["MonthLabel"] = df_year["MonthNum"].map(month_labels_map)

                # ---------------------------------------------------------
                # 📌 TOP MEMBERS WITH HIGHEST EXCUSED ABSENCES
                # ---------------------------------------------------------
                st.markdown("### 📝 Top Members With Highest **Excused Absences**")

                excused = df_year[df_year["Status"] == "Excused Absent"].copy()

                if excused.empty:
                    st.info("No excused absences this year.")
                else:
                    excused_summary = (
                        excused.groupby("Name")
                        .agg(
                            ExcusedCount=("Status", "size"),
                            ExcusedMonths=("MonthLabel", lambda x: ", ".join(sorted(set(x))))
                        )
                        .sort_values("ExcusedCount", ascending=False)
                        .reset_index()
                    )

                    st.dataframe(excused_summary, use_container_width=True, hide_index=True)

                    # Optional detailed insights
                    with st.expander("📌 Detailed View — Excused Absences"):
                        for _, row in excused_summary.iterrows():
                            st.markdown(
                                f"""
                                **{row['Name']}**  
                                - Excused Absences: **{row['ExcusedCount']}**  
                                - Months excused: **{row['ExcusedMonths']}**  
                                """
                            )

                # ---------------------------------------------------------
                # 📌 TOP MEMBERS WITH HIGHEST UNEXCUSED ABSENCES
                # ---------------------------------------------------------
                st.markdown("### ❌ Top Members With Highest **Unexcused Absences**")

                absent = df_year[df_year["Status"] == "Absent"].copy()

                if absent.empty:
                    st.info("No unexcused absences this year.")
                else:
                    absent_summary = (
                        absent.groupby("Name")
                        .agg(
                            AbsentCount=("Status", "size"),
                            AbsentMonths=("MonthLabel", lambda x: ", ".join(sorted(set(x))))
                        )
                        .sort_values("AbsentCount", ascending=False)
                        .reset_index()
                    )

                    st.dataframe(absent_summary, use_container_width=True, hide_index=True)

                    # Optional detailed insights
                    with st.expander("📌 Detailed View — Unexcused Absences"):
                        for _, row in absent_summary.iterrows():
                            st.markdown(
                                f"""
                                **{row['Name']}**  
                                - Unexcused Absences: **{row['AbsentCount']}**  
                                - Months absent: **{row['AbsentMonths']}**  
                                """
                            )


# ==========================================================================================================
# GENERATING REPORT PPT REPORT
# =============================================================================================================

            # ==========================================================================================================
            #                                         📊 POWERPOINT REPORT GENERATOR
            # ==========================================================================================================

            # ==========================================================================================================
            #                                         📊 POWERPOINT REPORT GENERATOR
            # ==========================================================================================================

            import io
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor

            st.divider()
            st.subheader("📥 Download Attendance PowerPoint Report")

            if st.button("📊 Download PPT Report"):
                prs = Presentation()

                # ================================
                # 1️⃣ TITLE SLIDE
                # ================================
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

                # Blue background rectangle
                bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                    prs.slide_width, prs.slide_height
                )
                bg.fill.solid()
                bg.fill.fore_color.rgb = RGBColor(30, 55, 110)
                bg.line.fill.background()

                # Title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.5))
                tf = title_box.text_frame
                tf.text = "OCA Attendance Report"
                tf.paragraphs[0].font.size = Pt(44)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Subtitle
                p = tf.add_paragraph()
                p.text = f"Year: {year_selected}"
                p.font.size = Pt(26)
                p.font.color.rgb = RGBColor(220, 220, 220)
                p.alignment = PP_ALIGN.CENTER

                # Logo
                logo_path = r"C:/Users/stans/OneDrive/Desktop/OCA/OCA LOGO/OCA LOGO.JPG"
                slide.shapes.add_picture(logo_path, Inches(3.4), Inches(4.6), width=Inches(3))

                # ================================
                # 2️⃣ PER-MEMBER STACKED BREAKDOWN
                # ================================
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.8))
                tbox.text_frame.text = "Per-Member Attendance Breakdown (Stacked)"
                tbox.text_frame.paragraphs[0].font.size = Pt(30)
                tbox.text_frame.paragraphs[0].font.bold = True

                img1 = io.BytesIO()
                fig1.savefig(img1, format="png", dpi=200, bbox_inches="tight")
                slide.shapes.add_picture(img1, Inches(0.4), Inches(1.0), width=Inches(9))

                # ================================
                # 3️⃣ MONTHLY ATTENDANCE BREAKDOWN
                # ================================
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.8))
                tbox.text_frame.text = "Monthly Attendance — Present, Absent, Excused Absent"
                tbox.text_frame.paragraphs[0].font.size = Pt(28)
                tbox.text_frame.paragraphs[0].font.bold = True

                img2 = io.BytesIO()
                fig3.savefig(img2, format="png", dpi=200, bbox_inches="tight")
                slide.shapes.add_picture(img2, Inches(0.4), Inches(1.0), width=Inches(9))

                # ================================
                # 4️⃣ TABLE — MEMBERS WITH 5+ MISSED
                # ================================
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.8))
                tbox.text_frame.text = "🚨 Members With 5+ Total Missed Meetings"
                tbox.text_frame.paragraphs[0].font.size = Pt(28)
                tbox.text_frame.paragraphs[0].font.bold = True

                rows = len(flagged) + 1
                cols = 4
                table = slide.shapes.add_table(
                    rows, cols, Inches(0.3), Inches(1.0), Inches(9), Inches(4)
                ).table

                headers = ["Name", "Total Missed", "Absent", "Excused"]
                for i, h in enumerate(headers):
                    table.cell(0, i).text = h

                for r, (_, row) in enumerate(flagged.iterrows(), start=1):
                    table.cell(r, 0).text = row["Name"]
                    table.cell(r, 1).text = str(row["TotalMissed"])
                    table.cell(r, 2).text = str(row["AbsentCount"])
                    table.cell(r, 3).text = str(row["ExcusedCount"])

                # ================================
                # 5️⃣ TOP EXCUSED ABSENCES (3+)
                # ================================
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.8))
                tbox.text_frame.text = "📝 Top Members With Highest Excused Absences"
                tbox.text_frame.paragraphs[0].font.size = Pt(28)
                tbox.text_frame.paragraphs[0].font.bold = True

                if not excused_summary.empty:
                    rows = len(excused_summary) + 1
                    table = slide.shapes.add_table(
                        rows, 3, Inches(0.3), Inches(1.0), Inches(9), Inches(4)
                    ).table

                    table.cell(0, 0).text = "Name"
                    table.cell(0, 1).text = "Excused Count"
                    table.cell(0, 2).text = "Months"

                    for r, (_, row) in enumerate(excused_summary.iterrows(), start=1):
                        table.cell(r, 0).text = row["Name"]
                        table.cell(r, 1).text = str(row["ExcusedCount"])
                        table.cell(r, 2).text = row["ExcusedMonths"]

                # ================================
                # 6️⃣ TOP UNEXCUSED ABSENCES (3+)
                # ================================
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.8))
                tbox.text_frame.text = "❌ Top Members With Highest Unexcused Absences"
                tbox.text_frame.paragraphs[0].font.size = Pt(28)
                tbox.text_frame.paragraphs[0].font.bold = True

                if not absent_summary.empty:
                    rows = len(absent_summary) + 1
                    table = slide.shapes.add_table(
                        rows, 3, Inches(0.3), Inches(1.0), Inches(9), Inches(4)
                    ).table

                    table.cell(0, 0).text = "Name"
                    table.cell(0, 1).text = "Absent Count"
                    table.cell(0, 2).text = "Months"

                    for r, (_, row) in enumerate(absent_summary.iterrows(), start=1):
                        table.cell(r, 0).text = row["Name"]
                        table.cell(r, 1).text = str(row["AbsentCount"])
                        table.cell(r, 2).text = row["AbsentMonths"]

                # ================================
                # EXPORT PPT
                # ================================
                out = io.BytesIO()
                prs.save(out)

                st.download_button(
                    "📥 Download OCA Attendance Report (PPTX)",
                    data=out.getvalue(),
                    file_name=f"OCA_Attendance_Report_{year_selected}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            # ============================================== ENDS HERE

        # ==========================================================================================================
        #                                         📄 POWERPOINT-STYLE PDF REPORT
        #      Full report: Title + Charts (fig1, fig3) + 5+ Missed + Excused + Unexcused Absence Tables
        # ==========================================================================================================

        import io
        from reportlab.platypus import (
            SimpleDocTemplate,
            Paragraph,
            Spacer,
            Image as RLImage,
            Table,
            TableStyle,
            PageBreak,
        )
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER
        from reportlab.lib import colors

        st.divider()
        st.subheader("📄 Download Attendance PDF Report")

        # Only show PDF button if we actually have analytics for the selected year
        if st.button("📄 Download PDF Report"):

            # -------------------------------
            # 1) Setup PDF document + styles
            # -------------------------------
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(
                buffer,
                pagesize=letter,
                title=f"OCA Attendance Report {year_selected}",
                leftMargin=40,
                rightMargin=40,
                topMargin=40,
                bottomMargin=40,
            )

            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                name="TitleCentered",
                parent=styles["Title"],
                alignment=TA_CENTER,
                fontSize=22,
                leading=26,
                textColor=colors.HexColor("#1e376e"),
            )
            subtitle_style = ParagraphStyle(
                name="SubtitleCentered",
                parent=styles["Normal"],
                alignment=TA_CENTER,
                fontSize=14,
                leading=18,
                textColor=colors.grey,
            )
            section_style = ParagraphStyle(
                name="SectionHeader",
                parent=styles["Heading2"],
                fontSize=16,
                leading=20,
                textColor=colors.HexColor("#1e376e"),
            )
            normal = styles["Normal"]

            story = []

            # -------------------------------
            # 2) TITLE PAGE WITH LOGO
            # -------------------------------
            logo_path = r"C:/Users/stans/OneDrive/Desktop/OCA/OCA LOGO/OCA LOGO.JPG"

            try:
                story.append(Spacer(1, 40))
                story.append(Paragraph("OCA Attendance Report", title_style))
                story.append(Spacer(1, 12))
                story.append(Paragraph(f"Year: {year_selected}", subtitle_style))
                story.append(Spacer(1, 24))
                story.append(Paragraph("Prepared by: OCA Publicity Secretary", normal))
                story.append(Spacer(1, 24))

                # Logo
                story.append(RLImage(logo_path, width=200, height=200))
            except Exception:
                # If logo missing, just skip it
                pass

            story.append(PageBreak())

            # -------------------------------
            # 3) PER-MEMBER STACKED CHART
            # -------------------------------
            try:
                story.append(Paragraph("Per-Member Attendance Breakdown (Stacked)", section_style))
                story.append(Spacer(1, 8))

                img_buf1 = io.BytesIO()
                fig1.savefig(img_buf1, format="png", dpi=150, bbox_inches="tight")
                img_buf1.seek(0)
                story.append(RLImage(img_buf1, width=500, height=260))
                story.append(Spacer(1, 16))
            except Exception:
                story.append(Paragraph("Per-member stacked chart unavailable.", normal))
                story.append(Spacer(1, 12))

            story.append(PageBreak())

            # -------------------------------
            # 4) MONTHLY ATTENDANCE CHART
            # -------------------------------
            try:
                story.append(Paragraph("Monthly Attendance — Present, Absent & Excused Absent", section_style))
                story.append(Spacer(1, 8))

                img_buf3 = io.BytesIO()
                fig3.savefig(img_buf3, format="png", dpi=150, bbox_inches="tight")
                img_buf3.seek(0)
                story.append(RLImage(img_buf3, width=500, height=260))
                story.append(Spacer(1, 16))
            except Exception:
                story.append(Paragraph("Monthly attendance chart unavailable.", normal))
                story.append(Spacer(1, 12))

            story.append(PageBreak())

            # -------------------------------
            # 5) TABLE — MEMBERS WITH 5+ MISSED
            # -------------------------------
            if "flagged" in locals() and not flagged.empty:
                story.append(Paragraph("🚨 Members With 5+ Total Missed Meetings", section_style))
                story.append(Spacer(1, 8))

                table_data = [["Name", "Total Missed", "Absent", "Excused", "Months"]]
                for _, row in flagged.iterrows():
                    table_data.append([
                        str(row["Name"]),
                        str(row["TotalMissed"]),
                        str(row["AbsentCount"]),
                        str(row["ExcusedCount"]),
                        str(row["MissedMonths"]),
                    ])

                t = Table(table_data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1e376e")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, 0), 10),
                    ("FONTSIZE", (0, 1), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.lightgrey]),
                ]))
                story.append(t)
                story.append(Spacer(1, 16))
            else:
                story.append(Paragraph("No member has 5 or more total missed meetings for this year.", normal))
                story.append(Spacer(1, 12))

            story.append(PageBreak())

            # -------------------------------
            # 6) TABLE — TOP EXCUSED ABSENCES
            # -------------------------------
            if "excused_summary" in locals() and not excused_summary.empty:
                story.append(Paragraph("📝 Members With Highest Excused Absences", section_style))
                story.append(Spacer(1, 8))

                table_data = [["Name", "Excused Count", "Months"]]
                for _, row in excused_summary.iterrows():
                    table_data.append([
                        str(row["Name"]),
                        str(row["ExcusedCount"]),
                        str(row["ExcusedMonths"]),
                    ])

                t = Table(table_data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2980b9")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, 0), 10),
                    ("FONTSIZE", (0, 1), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.lightgrey]),
                ]))
                story.append(t)
                story.append(Spacer(1, 16))
            else:
                story.append(Paragraph("No excused absences recorded for this year.", normal))
                story.append(Spacer(1, 12))

            story.append(PageBreak())

            # -------------------------------
            # 7) TABLE — TOP UNEXCUSED ABSENCES
            # -------------------------------
            if "absent_summary" in locals() and not absent_summary.empty:
                story.append(Paragraph("❌ Members With Highest Unexcused Absences", section_style))
                story.append(Spacer(1, 8))

                table_data = [["Name", "Absent Count", "Months"]]
                for _, row in absent_summary.iterrows():
                    table_data.append([
                        str(row["Name"]),
                        str(row["AbsentCount"]),
                        str(row["AbsentMonths"]),
                    ])

                t = Table(table_data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e74c3c")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, 0), 10),
                    ("FONTSIZE", (0, 1), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.lightgrey]),
                ]))
                story.append(t)
                story.append(Spacer(1, 16))
            else:
                story.append(Paragraph("No unexcused absences recorded for this year.", normal))
                story.append(Spacer(1, 12))

            # -------------------------------
            # 8) BUILD & DOWNLOAD
            # -------------------------------
            doc.build(story)
            buffer.seek(0)

            st.download_button(
                label="📥 Download OCA Attendance Report (PDF)",
                data=buffer.getvalue(),
                file_name=f"OCA_Attendance_Report_{year_selected}.pdf",
                mime="application/pdf",
            )

# # ============================= instruction  ====================
#
    # ---------------------------------------------------------
    # 📘 SIDEBAR TUTORIAL / USER GUIDE
    # ---------------------------------------------------------
    st.sidebar.markdown("---")   # <-- This adds a divider INSIDE SIDEBAR
    with st.sidebar.expander("📘 How to Use This System (Tutorial)", expanded=False):

        st.markdown("""
        ### 👋 Welcome to the OCA System Tutorial

        **This guide helps you navigate the Attendance & Financial Dashboard.**

        ---

        ## 🧍 1. Manage Members (Roster)
        - Add new members
        - Edit existing member names
        - Delete members and their attendance

        ---

        ## 📝 2. Take Attendance
        - Choose **Take New Attendance** to start a new record
        - Choose **Edit Existing Attendance** to update past records
        - Set default status and edit member notes

        ---

        ## 📄 3. Downloads
        - Export **CSV Roll-Call**, **Excel Roll-Call**, or **Template**

        ---

        ## ⬆️ 4. Upload Attendance
        - Import CSV/Excel files
        - System automatically updates the database

        ---

        ## 📊 5. Attendance Analytics
        - Filter by Year and Month
        - View charts and insights
        - Identify members with:
            - 5+ missed meetings
            - Highest excused absences
            - Highest unexcused absences

        ---

        ## 📘 6. PDF Report
        - Generate full attendance report
        - Includes charts, tables, insights

        ---

        ## 💰 7. OCA Financial Dashboard
        - View contributions
        - Track dues
        - Export presentations

        ---

        ## 🛠️ 8. Database
        - Database name: **oca_attendance.db**
        - Stores all members and attendance securely

        ---

        ### Need help?
        Contact ** Dr. Stanley Njoku**.
        """)


# ==============================================================================================================
# ATTENDANCE ENDS HERE
# ===============================================================================================================



#
# # =================================================================================== PART 33
#
# #
# def financial_ui():
#     import matplotlib.pyplot as plt
#     import matplotlib.ticker as mtick
#     from pptx import Presentation
#     from pptx.util import Inches, Pt
#     from pptx.enum.text import PP_ALIGN
#     from pptx.dml.color import RGBColor
#     import io
#     from datetime import datetime
#     import pandas as pd
#     import streamlit as st
    import os

    # ---------------- CSS for dashboard cards ----------------
    st.markdown("""
        <style>
        .cards {display:flex;flex-wrap:wrap;gap:18px;margin:15px 0 25px 0;}
        .card {flex:1 1 240px;background:linear-gradient(135deg,#f7f8fa,#eef2f7);
               border-radius:14px;padding:16px 18px;box-shadow:0 2px 10px rgba(0,0,0,0.08);
               border:1px solid #e4e6eb;}
        .label {font-size:13px;color:#555;text-transform:uppercase;letter-spacing:.5px;}
        .value {font-size:28px;font-weight:700;color:#222;margin-top:5px;}
        </style>
    """, unsafe_allow_html=True)

    st.title("💰 OCA Financial Data Analyzer")
    st.markdown("""
    Upload your **financial CSV file** (e.g., BMO bank statement) below.
    The system extracts *Posted Date, Description, Type, Credit/Debit, and Amount*
    and generates smart monthly summaries, annotated charts, and key highlights.
    """)

    uploaded_file = st.file_uploader("📂 Upload Financial CSV", type=["csv"])
    if uploaded_file is None:
        st.info("⬆️ Upload a CSV file to begin financial analysis.")
        return

    # ---------------- Read file ----------------
    try:
        df = pd.read_csv(uploaded_file)
    except UnicodeDecodeError:
        df = pd.read_csv(uploaded_file, encoding="ISO-8859-1")

    df.columns = df.columns.str.strip().str.upper()

    if "POSTED DATE" not in df.columns:
        st.error("❌ Missing 'POSTED DATE' column in file.")
        return

    df["POSTED DATE"] = pd.to_datetime(df["POSTED DATE"], errors="coerce")
    df["CREDIT/DEBIT"] = df["CREDIT/DEBIT"].astype(str).str.strip().str.title()
    df["MONTH"] = df["POSTED DATE"].dt.to_period("M")


    # ---------------- Monthly summary ----------------
    monthly_summary = df.groupby(["MONTH", "CREDIT/DEBIT"]).agg({"AMOUNT": "sum"}).reset_index()

    monthly_pivot = monthly_summary.pivot(
        index="MONTH",
        columns="CREDIT/DEBIT",
        values="AMOUNT"
    ).fillna(0)

    monthly_pivot["NET FLOW"] = monthly_pivot.get("Credit", 0) - monthly_pivot.get("Debit", 0)

    top_desc_type = df.groupby("MONTH").agg({
        "DESCRIPTION": lambda x: x.value_counts().index[0] if not x.empty else "",
        "TYPE": lambda x: x.value_counts().index[0] if not x.empty else ""
    })

    final_summary = monthly_pivot.join(top_desc_type).reset_index()
    final_summary["MONTH"] = pd.to_datetime(final_summary["MONTH"].astype(str)).dt.strftime("%b-%y")

    # ---------------- Dashboard cards ----------------
    total_credit = final_summary.get("Credit", pd.Series([0])).sum()
    total_debit = final_summary.get("Debit", pd.Series([0])).sum()
    net_balance = total_credit - total_debit
    avg_monthly_credit = final_summary.get("Credit", pd.Series([0])).mean()

    st.markdown(f"""
        <div class="cards">
          <div class="card"><div class="label">Total Credit</div><div class="value">${total_credit:,.2f}</div></div>
          <div class="card"><div class="label">Total Debit</div><div class="value">${total_debit:,.2f}</div></div>
          <div class="card"><div class="label">Net Balance</div><div class="value">${net_balance:,.2f}</div></div>
          <div class="card"><div class="label">Avg. Monthly Credit</div><div class="value">${avg_monthly_credit:,.2f}</div></div>
        </div>
    """, unsafe_allow_html=True)

    # ---------------- Summary table ----------------
    st.subheader("📅 Monthly Credit & Debit Summary")
    st.dataframe(final_summary, use_container_width=True, hide_index=True)

    # ---------------- Full dataset (Expandable) ----------------
    st.subheader("📂 Full Transaction Data")

    with st.expander(f"Click to view all transaction rows ({len(df)} records)", expanded=False):
        st.dataframe(df, use_container_width=True, hide_index=True)

        # Optional download button
        csv = df.to_csv(index=False).encode("utf-8")
        st.download_button(
            label="⬇️ Download Full Data",
            data=csv,
            file_name="transactions.csv",
            mime="text/csv"
        )

# # ======================================= ENDS HERE
#
#
#

    st.subheader("📅 Monthly Credit & Debit Summary")
    st.dataframe(final_summary, use_container_width=True, hide_index=True)

    # ---------------- Charts ----------------
    figs = {}
    st.divider()
    st.subheader("📈 Annotated Visual Trends")

    # 1️⃣ Monthly totals
    st.markdown("**Total Credit vs Debit (per Month)**")
    fig1, ax1 = plt.subplots(figsize=(8,4))
    months = final_summary["MONTH"]
    if "Credit" in final_summary:
        ax1.plot(months, final_summary["Credit"], marker="o", color="green", label="Credit")
        for i,v in enumerate(final_summary["Credit"]):
            ax1.annotate(f"${v:,.0f}", (i,v), textcoords="offset points", xytext=(0,8), ha="center", fontsize=8)
    if "Debit" in final_summary:
        ax1.plot(months, final_summary["Debit"], marker="o", color="red", label="Debit")
        for i,v in enumerate(final_summary["Debit"]):
            ax1.annotate(f"${v:,.0f}", (i,v), textcoords="offset points", xytext=(0,-12), ha="center", fontsize=8)
    ax1.set_ylabel("Amount ($)"); ax1.legend(); ax1.set_title("Monthly Totals")
    ax1.yaxis.set_major_formatter(mtick.StrMethodFormatter("${x:,.0f}"))
    st.pyplot(fig1); figs["Monthly Totals"] = fig1

    # 2️⃣ Net Flow
    st.markdown("**Net Flow (Credits – Debits)**")
    fig2, ax2 = plt.subplots(figsize=(8,4))
    bars = ax2.bar(final_summary["MONTH"], final_summary["NET FLOW"], color="#0078d4")
    for bar in bars:
        h = bar.get_height()
        ax2.annotate(f"${h:,.0f}", xy=(bar.get_x()+bar.get_width()/2, h), xytext=(0,5),
                     textcoords="offset points", ha="center", va="bottom", fontsize=8)
    ax2.set_title("Net Flow per Month"); ax2.set_ylabel("Net ($)")
    ax2.yaxis.set_major_formatter(mtick.StrMethodFormatter("${x:,.0f}"))
    st.pyplot(fig2); figs["Net Flow"] = fig2

    # ---------------- Single-click PPT Generation ----------------
    st.divider()
    st.subheader("📤 Export PowerPoint Report")

    if st.button("📥 Generate & Download OCA Financial Report (PPTX)"):
        prs = Presentation()
        slide_w, slide_h = prs.slide_width, prs.slide_height

        from pptx.dml.color import RGBColor
        import os, io

        # ---- Title Slide: Professional OCA Design ----
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background gradient
        bg_shape = title_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        fill = bg_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 51, 153)  # Deep blue
        bg_shape.line.fill.background()

        overlay = title_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 95, 200)
        overlay.fill.transparency = 0.25
        overlay.line.fill.background()

        # Add OCA logo
        logo_path = r"C:\Users\stans\OneDrive\Desktop\OCA\OCA LOGO\OCA LOGO.jpg"
        if os.path.exists(logo_path):
            title_slide.shapes.add_picture(logo_path, Inches(0.6), Inches(0.6), width=Inches(2))

        # --- Add professional finance-themed image ---
        # ✅ Updated image path
        finance_img_path = r"C:\Users\stans\OneDrive\Desktop\OCA\NEW OCA CHAPTER - 2025 - DR. ADMIN\Automated Tools\Images\Fin_image.jpg"
        if os.path.exists(finance_img_path):
            title_slide.shapes.add_picture(finance_img_path, Inches(6.6), Inches(0.6), width=Inches(3))
        else:
            from pptx.enum.shapes import MSO_SHAPE
            placeholder = title_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(0.6), Inches(3), Inches(2)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)
            placeholder.line.fill.background()
            placeholder.text = "Financial Image Placeholder"

        # Title
        title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(8), Inches(2))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = f"OCA FINANCIAL REPORT — {datetime.now():%B %Y}"
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Authors
        sub_box = title_slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(8), Inches(1.3))
        tf2 = sub_box.text_frame
        p2 = tf2.add_paragraph()
        p2.text = "Prepared by:"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(230, 230, 230)
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf2.add_paragraph()
        p3.text = "Ugo Harris — Financial Secretary\nJulie Ojibe — Assistant Financial Secretary"
        p3.font.size = Pt(20)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(255, 255, 255)
        p3.alignment = PP_ALIGN.CENTER

        # Footer
        footer_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        ft = footer_box.text_frame
        f = ft.add_paragraph()
        f.text = "Owerri Cultural Association (Northern California Chapter)"
        f.font.size = Pt(14)
        f.font.color.rgb = RGBColor(200, 200, 200)
        f.alignment = PP_ALIGN.CENTER

        # ---- Summary Slide ----
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tbox.text = "FINANCIAL SUMMARY"
        tf = tbox.text_frame
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True

        data = [
            ("Total Credit", f"${total_credit:,.2f}"),
            ("Total Debit", f"${total_debit:,.2f}"),
            ("Net Balance", f"${net_balance:,.2f}"),
            ("Avg Monthly Credit", f"${avg_monthly_credit:,.2f}")
        ]
        tbl = slide.shapes.add_table(rows=len(data), cols=2, left=Inches(1),
                                     top=Inches(1.5), width=Inches(8), height=Inches(2)).table
        for i,(lbl,val) in enumerate(data):
            tbl.cell(i,0).text = lbl
            tbl.cell(i,1).text = val

        # ---- Chart slides ----
        def add_slide(title, fig):
            buf = io.BytesIO()
            fig.savefig(buf, format="png", bbox_inches="tight", dpi=200)
            buf.seek(0)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8)).text = title
            slide.shapes.add_picture(buf, Inches(0.5), Inches(1), width=Inches(9))
            buf.close()

        for title, fig in figs.items():
            add_slide(title, fig)

        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        st.download_button(
            label="⬇️ Click to Download OCA Financial Report",
            data=pptx_io,
            file_name=f"OCA_Financial_Report_{datetime.now():%Y%m%d}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

# hjhjhjhjhjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjjj START HERE
        # ===================================================================================
        # 🔥 OCA DUES ANALYSIS (NEW SECTION - DOES NOT TOUCH FINANCIAL UI)

# PPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPPP  PART 2


#
def financial_ui():
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mtick
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import io
    from datetime import datetime
    import pandas as pd
    import streamlit as st

# def oca_dues_ui():

    import pandas as pd
    import streamlit as st
    import re
    from difflib import SequenceMatcher

    st.divider()
    st.subheader("👥 OCA Member Dues Analysis (Household Level)")

    uploaded_file = st.file_uploader(
        "📂 Upload Bank CSV for Dues Analysis",
        type=["csv"],
        key="dues_upload"
    )

    if uploaded_file is None:
        st.info("Upload bank statement to analyze dues")
        return

    # ===================================================================================
    # MASTER MEMBER LIST (FULL)
    # ===================================================================================
    members_data = [
        (1, "MR KELECHI ACHOLONU", 0), (1, "MRS NKIRU ACHOLONU", 0),
        (2, "MRS ROSEMARY ANYANWA", 0),
        (3, "MR HYACINTH I ANYASO CHIEF", 0), (3, "MRS LYDIA ANYASO LOLO", 0),
        (4, "DOC NNEKA CHUKWU", 0),
        (5, "MR CHIAGORO CHUKWUMA", 0),
        (6, "MR EMMA DIALA", 0), (6, "MRS JOYCE DIALA", 0),
        (7, "MRS GLORIA DURU", 0),
        (8, "MRS LUCY ECHIBE", 0),
        (9, "MRS AGNES EGU", 0),
        (10, "MRS ROSE EGU", 0),
        (11, "MR ALEXANDER A EKE", 0), (11, "MRS LINDA EKE", 35.01),
        (12, "MR DAMIAN EKE CHIEF", 0), (12, "MR ULOMA EKE LOLO", 0),
        (13, "MR CYPRIAN EGEOLU", 0), (13, "MRS EGEOLU", 0),
        (14, "MRS ADAMNA EMEZIE", 0),
        (15, "MR PRINCE IBE", 0), (15, "MRS CHINWE IBE", 0),
        (16, "MR ANTHONY IHEJITO", 0),
        (17, "DOC STANLEY NJOKU", 0), (17, "MRS AMAKA NJOKU", 0),
        (18, "REV SISTER AKUNNA NJOKU", 0),
        (19, "MR LINUS NWAULU", 0), (19, "MRS NWAULU", 0),
        (20, "MR PETER OBILOR", 0),
        (21, "MRS VIVIAN OBICHERE", 0),
        (22, "MR ETHEBERT OGBUEHI", 0),
        (23, "MR KINSLEY OGWUDIRE", 0), (23, "MRS IJEOMA OGWUDIRE", 0),
        (24, "MRS JULIANA OJIBE", 0),
        (25, "MR THEOPHILUS N ONYENEKE", 0), (25, "MRS BEATRICE ONYENEKE", 35.01),
        (26, "SIR ETHELBERT R ONYEWUNYI", 0),
        (27, "MR SAMUEL O ONYENWEE", 0), (27, "MRS ONYENWE", 0),
        (28, "MR DAMIAN ONYEUKWU", 0),
        (29, "MR EMEKA OPARAOCHEKWE", 0),
        (30, "MR WILSON UDENJI", 0),
        (31, "MR OBI UGORJI", 0), (31, "MRS CHINONYE UGORJI", 0),
        (32, "DOC GIBSON UNAJI", 0),
        (33, "MR AUSTIN UWAKWE", 0), (33, "MRS PAMELA UWAKWE", 0),
        (34, "MR DAVID UZOMA UZOHUO", 0), (34, "MRS UZOHUO", 0),
        (35, "MRS JOSEPHINE YOKO UZOMA", 0),

        (36, "MRS BEATRICE EGU", 17.21),
        (37, "MR ERNEST EGU", 34.41),
        (38, "MRS JOY EGU", 17.21),
        (39, "MRS LILIAN EGU", 34.41), (39, "MR ROBERT EGU", 17.21),
        (40, "MRS BEATRICE EKE", 34.41),
        (41, "MRS AUGUSTINA EKEH", 17.81), (41, "MR GEORGE EKEH", 17.21),
        (42, "MISS LORETTA OGECHI IWU", 34.41),
        (43, "MRS JACINTA MBATA", 35.01), (43, "MR CHRISTOPHER ORJI", 17.21),
        (44, "MRS MARGARET P NWACHUKWU AKAW", 17.21),
        (45, "LADY JOANNES NWADIBIA", 17.81),
        (46, "MRS CONSTANCE NWADIKE", 35.01),
        (47, "MR AUGUSTINE NWAOGU", 35.01), (47, "MRS NGOZI NWAOGU", 35.01),
        (48, "MR ADOLF OBILOR", 17.21), (48, "MRS MARGARET NGOZI OBILOR", 35.01),
        (49, "MR INNOCENT OBILOR", 17.21),
        (50, "MR CASMIR OBINNA", 35.01),
        (51, "MRS ELIZABERTH OHA", 34.41), (51, "MR LINUS OHA", 35.01),
        (52, "MR PAUL OHA", 35.01), (52, "MRS JOYCE OHA", 35.01),
        (53, "MRS ROSA ONYEAGUCHA", 34.41),
        (54, "MRS FIDELIA ONYEJIKWE", 17.21),
        (55, "MRS GLORIA OPARA", 34.41), (55, "MR UGO HARRIS OPARA", 34.41),
        (56, "MR IK F ORJI", 17.21), (56, "MRS VERONICA ORJI", 35.01),
        (57, "MRS CELINE OSUALA", 17.21),
        (58, "MR AUGUSTINE OWOBETE", 34.41), (58, "MRS GRACE N OWOBETE", 35.01),
        (59, "MR CHINONYE UZOMA", 35.01), (59, "MRS MARINA UZOMA", 35.01),
        (60, "MR TOM MBARA", 0), (60, "MRS EMELDA MBARA", 0),
    ]


# ]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]] part

    members_df = pd.DataFrame(members_data, columns=["COUPLE_ID", "NAME", "INSURANCE"])

    BASE_DUE = 10
    grouped = members_df.groupby("COUPLE_ID")

    couple_df = grouped.agg({
        "NAME": lambda x: " & ".join(x),
        "INSURANCE": "sum"
    }).reset_index()

    couple_df["DUES_TOTAL"] = grouped.size().values * BASE_DUE
    couple_df.rename(columns={"INSURANCE": "INSURANCE_TOTAL"}, inplace=True)
    couple_df["MONTHLY_TOTAL"] = couple_df["DUES_TOTAL"] + couple_df["INSURANCE_TOTAL"]

    # ===================================================================================
    # LOAD BANK
    # ===================================================================================
    df = pd.read_csv(uploaded_file, encoding="ISO-8859-1")
    df.columns = df.columns.str.upper().str.strip()

    df["POSTED DATE"] = pd.to_datetime(df["POSTED DATE"], errors="coerce")
    df["DESCRIPTION"] = df["DESCRIPTION"].astype(str).str.upper()
    df["AMOUNT"] = pd.to_numeric(df["AMOUNT"], errors="coerce").fillna(0)
    df["CREDIT/DEBIT"] = df["CREDIT/DEBIT"].str.title()

    payments = df[df["CREDIT/DEBIT"] == "Credit"].copy()

    # ===================================================================================
    # CHECK INPUT
    # ===================================================================================
    st.subheader("➕ Add Check Payment")

    if "manual_checks" not in st.session_state:
        st.session_state.manual_checks = []

    name_input = st.selectbox("Select Member", members_df["NAME"].unique())
    amt_input = st.number_input("Amount", min_value=0.0, step=10.0)
    date_input = st.date_input("Date")

    if st.button("Add Check"):
        st.session_state.manual_checks.append({
            "NAME": name_input,
            "AMOUNT": amt_input,
            "POSTED DATE": pd.to_datetime(date_input),
            "CREDIT/DEBIT": "Credit"
        })
        st.success("Check added")

    if st.session_state.manual_checks:
        st.dataframe(pd.DataFrame(st.session_state.manual_checks))

    # ===================================================================================
    # MATCHING
    # ===================================================================================
    def clean_tokens(x):
        x = re.sub(r"[^A-Z\s]", " ", str(x).upper())
        remove = ["MR", "MRS", "DOC", "DR", "MISS", "SIR", "REV", "SISTER", "FROM", "ZELLE"]
        for r in remove:
            x = x.replace(r, " ")
        return set(x.split())

    ALIAS_MAP = {
        "THOMAS MBARA": "TOM MBARA",
        "ROY ONYEWUNYI": "ETHELBERT R ONYEWUENYI"
    }

    def match_name(raw):
        raw = raw.upper()
        for a, v in ALIAS_MAP.items():
            if a in raw:
                raw = v

        rt = clean_tokens(raw)
        best = None
        best_score = 0

        for _, row in members_df.iterrows():
            mt = clean_tokens(row["NAME"])
            overlap = len(rt & mt) / max(len(mt), 1)
            char = SequenceMatcher(None, " ".join(rt), " ".join(mt)).ratio()
            score = max(overlap, char)

            if score > best_score:
                best_score = score
                best = row["NAME"]

        return best if best_score >= 0.65 else None

    payments["NAME_RAW"] = payments["DESCRIPTION"]
    payments["NAME"] = payments["DESCRIPTION"].apply(match_name)

    # ===================================================================================
    # MERGE CHECKS
    # ===================================================================================
    if st.session_state.manual_checks:
        check_df = pd.DataFrame(st.session_state.manual_checks)
        payments = pd.concat([payments, check_df], ignore_index=True)

    payments["COUPLE_ID"] = payments["NAME"].map(
        members_df.set_index("NAME")["COUPLE_ID"]
    )

    matched = payments[payments["COUPLE_ID"].notna()]
    unmatched = payments[payments["COUPLE_ID"].isna()]

    # ===================================================================================
    # AGG
    # ===================================================================================
    paid = matched.groupby("COUPLE_ID")["AMOUNT"].sum().reset_index()
    paid.rename(columns={"AMOUNT": "TOTAL_PAID"}, inplace=True)

    report = couple_df.merge(paid, on="COUPLE_ID", how="left")
    report["TOTAL_PAID"] = report["TOTAL_PAID"].fillna(0)

    # ===================================================================================
    # CALC
    # ===================================================================================
    CURRENT_MONTH = pd.Timestamp.today().to_period("M")
    START_MONTH = pd.Period("2026-01")

    months_due = (CURRENT_MONTH - START_MONTH).n + 1

    report["EXPECTED"] = report["MONTHLY_TOTAL"] * months_due

    raw_balance = report["EXPECTED"] - report["TOTAL_PAID"]

    report["BALANCE"] = raw_balance.clip(lower=0)
    report["PAID_AHEAD_BY"] = raw_balance.apply(lambda x: abs(x) if x < 0 else 0)

    report["MONTHS_BEHIND"] = (report["BALANCE"] / report["MONTHLY_TOTAL"]).astype(int)

    # ======================================= MODIFY
    report["MONTHS_BEHIND"] = (report["BALANCE"] / report["MONTHLY_TOTAL"]).astype(int)

    # ===================================================================================
    # 🔥 STATUS COLUMN
    # ===================================================================================
    report["MONTHS_PAID"] = (report["TOTAL_PAID"] / report["MONTHLY_TOTAL"]).fillna(0).astype(int)

    def get_status(row):
        if row["MONTHS_PAID"] >= 12:
            return "PAID IN FULL"
        elif row["MONTHS_BEHIND"] <= 3:
            return "CURRENT"
        else:
            return "BEHIND"

    report["STATUS"] = report.apply(get_status, axis=1)


    # ========================================================================= FILTER MEMBERS PAID IN FULL AND
    # ========================================================================= FILTER MEMBERS PAID IN FULL AND
    # ========================================================================= FILTER MEMBERS PAID IN FULL AND
    # ========================================================================= FILTER MEMBERS PAID IN FULL AND


    # ===================================================================================
    # DISPLAY
    # ===================================================================================
    st.dataframe(report, use_container_width=True)

    # ===================================================================================
    # 📊 STATUS TABLES (PAID / CURRENT / BEHIND)
    # ===================================================================================
    st.subheader("📊 Member Status Breakdown")

    # 🟢 PAID IN FULL
    paid_full = report[report["STATUS"] == "PAID IN FULL"].copy()

    st.markdown(f"### 🟢 Paid In Full ({len(paid_full)})")
    st.dataframe(
        paid_full[
            [
                "COUPLE_ID",
                "NAME",
                "MONTHLY_TOTAL",
                "TOTAL_PAID",
                "EXPECTED",
                "BALANCE",
                "MONTHS_PAID",
                "STATUS"
            ]
        ],
        use_container_width=True
    )

    # 🟡 CURRENT
    current = report[report["STATUS"] == "CURRENT"].copy()

    st.markdown(f"### 🟡 Current (≤ 3 Months Behind) ({len(current)})")
    st.dataframe(
        current[
            [
                "COUPLE_ID",
                "NAME",
                "MONTHLY_TOTAL",
                "TOTAL_PAID",
                "EXPECTED",
                "BALANCE",
                "MONTHS_BEHIND",
                "STATUS"
            ]
        ],
        use_container_width=True
    )

    # 🔴 BEHIND (3+ MONTHS)
    behind = report[report["MONTHS_BEHIND"] >= 3].copy()

    # Optional sort (most owed first)
    behind = behind.sort_values("BALANCE", ascending=False)

    st.markdown(f"### 🔴 Behind (3+ Months) ({len(behind)})")
    st.dataframe(
        behind[
            [
                "COUPLE_ID",
                "NAME",
                "MONTHLY_TOTAL",
                "TOTAL_PAID",
                "EXPECTED",
                "BALANCE",
                "MONTHS_BEHIND",
                "STATUS"
            ]
        ],
        use_container_width=True
    )
    #



    # 📊 ENHANCED SUMMARY (WITH TOTAL AMOUNT OWED)
    # ===================================================================================
    st.subheader("📊 Summary")

    summary = report.groupby("MONTHS_BEHIND").agg(
        COUNT=("COUPLE_ID", "count"),
        TOTAL_BALANCE=("BALANCE", "sum")
    ).reset_index().sort_values("MONTHS_BEHIND", ascending=False)

    # 🔥 Format money
    summary["TOTAL_BALANCE"] = summary["TOTAL_BALANCE"].apply(lambda x: f"${x:,.2f}")

    st.dataframe(summary, use_container_width=True)

    with st.expander(f"⚠️ Unmatched Payments ({len(unmatched)})"):
        st.dataframe(unmatched[["POSTED DATE", "DESCRIPTION", "AMOUNT", "NAME_RAW"]])



# ]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]] DISPLAY THOSE CURRENT










# ]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]
# SUMMARY AND GRAPHS FOR MONEY RECIEVED
# \\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\
    # ===================================================================================
    # 💰 PAYMENT RECONCILIATION SUMMARY
    # ===================================================================================

    st.divider()

    st.subheader("💰 Payment Reconciliation Summary")

    # 🔥 Total received (ALL money: Zelle + Checks)
    total_received = payments["AMOUNT"].sum()

    # 🔥 Assigned (matched to households)
    total_assigned = matched["AMOUNT"].sum()

    # 🔥 Unassigned (not matched)
    total_unassigned = unmatched["AMOUNT"].sum()

    # 🔥 Sanity check
    difference = total_received - (total_assigned + total_unassigned)

    col1, col2, col3 = st.columns(3)

    with col1:
        st.metric("Total Received", f"${total_received:,.2f}")

    with col2:
        st.metric("Assigned to Households", f"${total_assigned:,.2f}")

    with col3:
        st.metric("Unassigned", f"${total_unassigned:,.2f}")

    # 🔥 Validation check
    if abs(difference) > 0.01:
        st.error(f"⚠️ Reconciliation mismatch: ${difference:,.2f}")
    else:
        st.success("✅ Reconciliation balanced")

    # ===================================================================================
    # 🔍 OPTIONAL: SHOW UNASSIGNED DETAILS
    # ===================================================================================
    with st.expander(f"🔍 View Unassigned Payments (${total_unassigned:,.2f})"):
        st.dataframe(
            unmatched[["POSTED DATE", "DESCRIPTION", "AMOUNT", "NAME_RAW"]],
            use_container_width=True
        )


# gggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggggg START HERE
    # ===================================================================================
    # 📈 MONTHLY CASHFLOW ANALYSIS (FINAL VERSION - PPT READY)
    # ===================================================================================
    import matplotlib.pyplot as plt

    st.subheader("📈 Monthly Cashflow Analysis")

    # -----------------------------------------------------------------------------------
    # PREP DATA
    # -----------------------------------------------------------------------------------
    df["POSTED DATE"] = pd.to_datetime(df["POSTED DATE"], errors="coerce")
    df = df[df["POSTED DATE"].notna()]

    df["MONTH"] = df["POSTED DATE"].dt.to_period("M").astype(str)

    # 💰 Money In
    monthly_in = df[df["CREDIT/DEBIT"] == "Credit"].groupby("MONTH")["AMOUNT"].sum()

    # 💸 Money Out (FORCE POSITIVE)
    monthly_out = df[df["CREDIT/DEBIT"] == "Debit"].groupby("MONTH")["AMOUNT"].sum().abs()

    # Combine
    monthly_df = pd.DataFrame({
        "Money In ($)": monthly_in,
        "Money Out ($)": monthly_out
    }).fillna(0).reset_index()

    # -----------------------------------------------------------------------------------
    # SORT MONTH
    # -----------------------------------------------------------------------------------
    monthly_df["MONTH"] = pd.to_datetime(monthly_df["MONTH"])
    monthly_df = monthly_df.sort_values("MONTH")
    monthly_df["MONTH_STR"] = monthly_df["MONTH"].dt.strftime("%b-%Y")

    # -----------------------------------------------------------------------------------
    # 🔥 DIFF (CORRECT LOGIC)
    # -----------------------------------------------------------------------------------
    monthly_df["Diff ($)"] = monthly_df["Money In ($)"] - monthly_df["Money Out ($)"]

    # -----------------------------------------------------------------------------------
    # 📊 FORMAT TABLE ($ + ,)
    # -----------------------------------------------------------------------------------
    display_df = monthly_df.copy()

    display_df["Money In ($)"] = display_df["Money In ($)"].apply(lambda x: f"${x:,.2f}")
    display_df["Money Out ($)"] = display_df["Money Out ($)"].apply(lambda x: f"${x:,.2f}")
    display_df["Diff ($)"] = display_df["Diff ($)"].apply(lambda x: f"${x:,.2f}")

    st.dataframe(
        display_df[["MONTH_STR", "Money In ($)", "Money Out ($)", "Diff ($)"]],
        use_container_width=True
    )

    # ===================================================================================
    # 📈 LINE CHART (ANNOTATED)
    # ===================================================================================
    st.subheader("📊 Monthly Trend (Annotated)")

    fig, ax = plt.subplots(figsize=(10, 5))

    ax.plot(monthly_df["MONTH_STR"], monthly_df["Money In ($)"], marker='o')
    ax.plot(monthly_df["MONTH_STR"], monthly_df["Money Out ($)"], marker='o')

    # Annotate values
    for i, row in monthly_df.iterrows():
        ax.text(i, row["Money In ($)"], f'${row["Money In ($)"]:,.0f}', ha='center', va='bottom')
        ax.text(i, row["Money Out ($)"], f'${row["Money Out ($)"]:,.0f}', ha='center', va='top')

    ax.set_title("Monthly Cashflow")
    ax.set_xlabel("Month")
    ax.set_ylabel("Amount ($)")
    ax.legend(["Money In", "Money Out"])

    plt.xticks(rotation=69)  # ✅ ROTATE LABELS

    st.pyplot(fig)

    # ===================================================================================
    # 📊 BAR CHART (ANNOTATED)
    # ===================================================================================
    st.subheader("📊 Monthly Comparison")

    fig_bar, ax_bar = plt.subplots(figsize=(10, 5))

    x = range(len(monthly_df))

    bars1 = ax_bar.bar(x, monthly_df["Money In ($)"])
    bars2 = ax_bar.bar(x, monthly_df["Money Out ($)"])

    ax_bar.set_xticks(x)
    ax_bar.set_xticklabels(monthly_df["MONTH_STR"])

    # Annotate bars
    for bar in bars1:
        height = bar.get_height()
        ax_bar.text(bar.get_x() + bar.get_width() / 2, height,
                    f'${height:,.0f}', ha='center', va='bottom')

    for bar in bars2:
        height = bar.get_height()
        ax_bar.text(bar.get_x() + bar.get_width() / 2, height,
                    f'${height:,.0f}', ha='center', va='top')

    ax_bar.set_title("Monthly Comparison")
    ax_bar.set_ylabel("Amount ($)")
    ax_bar.legend(["Money In", "Money Out"])

    plt.xticks(rotation=69)  # ✅ ROTATE LABELS

    st.pyplot(fig_bar)

    # ===================================================================================
    # 💰 DIFF GRAPH (FINAL)
    # ===================================================================================
    st.subheader("💰 Monthly Difference (Diff)")

    fig_diff, ax_diff = plt.subplots(figsize=(10, 5))

    ax_diff.plot(monthly_df["MONTH_STR"], monthly_df["Diff ($)"], marker='o')

    # Annotate Diff values
    for i, row in monthly_df.iterrows():
        ax_diff.text(i, row["Diff ($)"], f'${row["Diff ($)"]:,.0f}',
                     ha='center', va='bottom')

    ax_diff.set_title("Monthly Difference")
    ax_diff.set_xlabel("Month")
    ax_diff.set_ylabel("Amount ($)")
    ax_diff.legend(["Diff"])

    plt.xticks(rotation=69)  # ✅ ROTATE LABELS

    st.pyplot(fig_diff)

    st.divider()

# ]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]
    # ===================================================================================
    # 📊 COLLECTION PERFORMANCE (YEAR vs ACTUAL)
    # ===================================================================================
    st.subheader("📊 Collection Performance Overview")

    # -----------------------------------------------------------------------------------
    # 🔥 TIME LOGIC
    # -----------------------------------------------------------------------------------
    current_period = pd.Timestamp.today().to_period("M")
    start_period = pd.Period("2026-01")
    end_period = pd.Period("2026-12")

    months_elapsed = (current_period - start_period).n + 1
    total_year_months = 12

    # -----------------------------------------------------------------------------------
    # 💰 EXPECTED CALCULATIONS
    # -----------------------------------------------------------------------------------
    # Total monthly obligation across all households
    total_monthly_expected = report["MONTHLY_TOTAL"].sum()

    # Expected by now
    expected_to_date = total_monthly_expected * months_elapsed

    # Expected full year
    expected_full_year = total_monthly_expected * total_year_months

    # -----------------------------------------------------------------------------------
    # 💵 ACTUAL RECEIVED
    # -----------------------------------------------------------------------------------
    # total_received = report["TOTAL_PAID"].sum()

    total_received = payments["AMOUNT"].sum()

    # -----------------------------------------------------------------------------------
    # 📊 VARIANCE
    # -----------------------------------------------------------------------------------
    variance_to_date = total_received - expected_to_date
    variance_full_year = total_received - expected_full_year

    # -----------------------------------------------------------------------------------
    # 📊 DISPLAY METRICS
    # -----------------------------------------------------------------------------------
    col1, col2 = st.columns(2)

    with col1:
        st.markdown("### 📅 Year-to-Date Performance")
        st.metric(
            "Expected (YTD)",
            f"${expected_to_date:,.2f}"
        )
        st.metric(
            "Received (YTD)",
            f"${total_received:,.2f}"
        )
        st.metric(
            "Difference",
            f"${variance_to_date:,.2f}"
        )

    with col2:
        st.markdown("### 📆 Full Year Projection")
        st.metric(
            "Expected (Full Year)",
            f"${expected_full_year:,.2f}"
        )
        st.metric(
            "Received So Far",
            f"${total_received:,.2f}"
        )
        st.metric(
            "Gap to Year Target",
            f"${variance_full_year:,.2f}"
        )

    # -----------------------------------------------------------------------------------
    # 🔥 COLLECTION RATE
    # -----------------------------------------------------------------------------------
    collection_rate = (total_received / expected_to_date * 100) if expected_to_date > 0 else 0

    st.metric("📈 Collection Rate (YTD)", f"{collection_rate:.1f}%")

    # -----------------------------------------------------------------------------------
    # 🚨 STATUS MESSAGE
    # -----------------------------------------------------------------------------------
    if variance_to_date >= 0:
        st.success("✅ Collections are on track or ahead")
    else:
        st.warning("⚠️ Collections are behind schedule")

    st.divider()

    # ===================================================================================
    # 📊 COLLECTION PERFORMANCE GRAPH (PPT READY)
    # ===================================================================================
    import matplotlib.pyplot as plt

    st.subheader("📊 Collection Performance (Visual)")

    labels = ["Expected (YTD)", "Received (YTD)", "Gap"]

    values = [
        expected_to_date,
        total_received,
        variance_to_date
    ]

    fig, ax = plt.subplots(figsize=(8, 5))

    bars = ax.bar(labels, values)

    # Annotate values ($ + comma)
    for bar in bars:
        height = bar.get_height()
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            height,
            f'${height:,.0f}',
            ha='center',
            va='bottom'
        )

    # Titles
    ax.set_title("Expected vs Actual Collection")
    ax.set_ylabel("Amount ($)")

    st.pyplot(fig)

    # ===================================================================================
    # 📊 ANNUAL COLLECTION PERFORMANCE GRAPH
    # ===================================================================================
    st.subheader("📆 Annual Collection Benchmark")

    import matplotlib.pyplot as plt

    # 🔥 VALUES
    annual_labels = ["Expected (Full Year)", "Received (So Far)", "Remaining Gap"]

    remaining_gap = expected_full_year - total_received

    annual_values = [
        expected_full_year,
        total_received,
        remaining_gap
    ]

    fig2, ax2 = plt.subplots(figsize=(8, 5))

    bars = ax2.bar(annual_labels, annual_values)

    # 🔥 Annotate values ($ + commas)
    for bar in bars:
        height = bar.get_height()
        ax2.text(
            bar.get_x() + bar.get_width() / 2,
            height,
            f'${height:,.0f}',
            ha='center',
            va='bottom'
        )

    # Titles
    ax2.set_title("Annual Expected vs Actual Collection")
    ax2.set_ylabel("Amount ($)")

    st.pyplot(fig2)



# \\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\ POWER POINT slides
#

# =========================================================================PAER 2

    # ===================================================================================
    # 📊 POWERPOINT FINANCIAL REPORT GENERATOR (FINAL - NO OVERFLOW)
    # ===================================================================================
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.dml.color import RGBColor
    import matplotlib.pyplot as plt
    from datetime import datetime
    import os

    st.divider()
    st.subheader("📥 Download Financial PowerPoint Report")

    if st.button("📊 Download Financial PPT Report"):

        prs = Presentation()

        # ===================================================================================
        # 1️⃣ TITLE SLIDE (UNCHANGED DESIGN)
        # ===================================================================================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        
        BASE_DIR = os.path.dirname(__file__)
        logo_path = os.path.join(BASE_DIR, "assets", "FNEW_ACE_LOGO.png")     # NEW LOGO UPDATE 4/10/2026


        

        BASE_DIR = os.path.dirname(__file__)
        logo_path = os.path.join(BASE_DIR, "assets", "FNEW_ACE_LOGO.png")
        
        print("DEBUG PATH:", logo_path)
        print("FILE EXISTS:", os.path.exists(logo_path))
        
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(3.5), Inches(0.3), width=Inches(2))
        else:
            print("❌ Logo NOT FOUND")







        

        #logo_path = r"C:\Users\stans\OneDrive\Desktop\OCA\01 -STANLEY'S ADMINISTRATION - PRESIDENT\OCA NEW LOGO\OCA - FNEW ACE LOGO.png"

        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(3.5), Inches(0.3), width=Inches(2))

        org_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        tf = org_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Owerri Cultural Association (OCA)"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 128)
        p.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Financial Report"
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        current_month = datetime.now().strftime("%B %Y")

        p2 = tf.add_paragraph()
        p2.text = current_month
        p2.font.size = Pt(24)
        p2.alignment = PP_ALIGN.CENTER

        info_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(2))
        tf = info_box.text_frame

        tf.text = "Financial Secretary: Mr. Ugo Harris Opara"
        tf.add_paragraph().text = "Assistant Financial Secretary: Mrs. Pamela Uwakwe"
        tf.add_paragraph().text = "Treasurer: Lady Ogechi Iwu"

        footer = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        footer.text_frame.text = f"Generated on {datetime.now():%B %d, %Y}"

        # ===================================================================================
        # HELPER: ADD CHART SLIDE (UNCHANGED)
        # ===================================================================================
        def add_slide(title, fig):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            tf = tbox.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.bold = True

            img = io.BytesIO()
            fig.savefig(img, format="png", dpi=200, bbox_inches="tight")
            img.seek(0)

            slide.shapes.add_picture(img, Inches(0.5), Inches(1.2), width=Inches(9))

        # ===================================================================================
        # HELPER: TABLE (8 ROWS PER SLIDE - FINAL FIX)
        # ===================================================================================
        def add_table_slides(title, df, columns, rows_per_slide=8):

            total_rows = len(df)
            num_slides = (total_rows // rows_per_slide) + 1

            for s in range(num_slides):

                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # Title
                tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                tf = tbox.text_frame
                tf.text = f"{title} (Part {s + 1})"
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True

                chunk = df.iloc[s * rows_per_slide:(s + 1) * rows_per_slide]

                rows = len(chunk) + 1
                cols = len(columns)

                table = slide.shapes.add_table(
                    rows, cols, Inches(0.3), Inches(1.2), Inches(9), Inches(4.5)
                ).table

                # 🔥 COLUMN WIDTH FIX
                table.columns[0].width = Inches(4.0)
                for i in range(1, cols):
                    table.columns[i].width = Inches(1.6)

                # Headers
                for i, col in enumerate(columns):
                    cell = table.cell(0, i)
                    cell.text = col
                    cell.text_frame.paragraphs[0].font.bold = True

                # Data
                for r, (_, row) in enumerate(chunk.iterrows(), start=1):
                    for c, col in enumerate(columns):

                        val = row[col]

                        if isinstance(val, (int, float)):
                            val = f"${val:,.0f}"

                        cell = table.cell(r, c)
                        tf = cell.text_frame
                        tf.clear()

                        p = tf.paragraphs[0]
                        p.text = str(val)

                        tf.word_wrap = True
                        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                # 🔥 ROW HEIGHT FIX
                for i in range(rows):
                    table.rows[i].height = Inches(0.6)

        # ===================================================================================
        # CHARTS (UNCHANGED)
        # ===================================================================================
        add_slide("Monthly Cashflow Trend", fig)
        add_slide("Monthly Comparison", fig_bar)
        add_slide("Monthly Difference", fig_diff)

        # ===================================================================================
        # 📊 MONTHLY CASHFLOW TABLE
        # ===================================================================================
        cashflow_table = monthly_df.rename(columns={
            "MONTH_STR": "Month",
            "Money In ($)": "Money In",
            "Money Out ($)": "Money Out",
            "Diff ($)": "Diff"
        })

        add_table_slides(
            "Monthly Cashflow Table",
            cashflow_table,
            ["Month", "Money In", "Money Out", "Diff"]
        )

        # ===================================================================================
        # 🚨 DELINQUENCY TABLE
        # ===================================================================================
        clean_high_due = report[report["MONTHS_BEHIND"] >= 3].copy()

        add_table_slides(
            "🚨 Households Owing 3+ Months",
            clean_high_due,
            ["NAME", "TOTAL_PAID", "EXPECTED", "BALANCE"]
        )

        # ===================================================================================
        # ===================================================================================
        # 📊 SUMMARY TABLE (POWERPOINT SLIDE) - FIXED
        # ===================================================================================
        def add_summary_slide(summary_df):

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Title
            tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            tf = tbox.text_frame
            tf.text = "📊 Summary"
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.bold = True

            rows = len(summary_df) + 1
            cols = 3

            table = slide.shapes.add_table(
                rows, cols, Inches(1), Inches(1.5), Inches(7), Inches(4)
            ).table

            # Column headers
            headers = ["Months Behind", "Count", "Total Balance"]
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = h
                cell.text_frame.paragraphs[0].font.bold = True

            # Data rows (SAFE FIX APPLIED)
            for r, (_, row) in enumerate(summary_df.iterrows(), start=1):

                # Months Behind
                table.cell(r, 0).text = str(int(row["MONTHS_BEHIND"]))

                # Count
                table.cell(r, 1).text = str(int(row["COUNT"]))

                # 🔥 SAFE MONEY FORMAT (NO MORE ERROR)
                val = row["TOTAL_BALANCE"]
                val = str(val).replace("$", "").replace(",", "")

                try:
                    val = float(val)
                    formatted = f"${val:,.2f}"
                except:
                    formatted = str(val)

                table.cell(r, 2).text = formatted

            # Column widths
            table.columns[0].width = Inches(2.5)
            table.columns[1].width = Inches(2)
            table.columns[2].width = Inches(2.5)

        # ===================================================================================
        # CALL FUNCTION
        # ===================================================================================
        add_summary_slide(summary)

        # ===================================================================================
        # 🙏 THANK YOU SLIDE (FINAL PAGE)
        # ===================================================================================
        def add_thank_you_slide():

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Main message
            tbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
            tf = tbox.text_frame

            p = tf.paragraphs[0]
            p.text = "🙏 Thank You for Your Time"
            p.font.size = Pt(40)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Optional subtitle
            p2 = tf.add_paragraph()
            p2.text = "We appreciate your attention and support."
            p2.font.size = Pt(20)
            p2.alignment = PP_ALIGN.CENTER

        # ===================================================================================
        # CALL THANK YOU SLIDE (LAST SLIDE)
        # ===================================================================================
        add_thank_you_slide()























        # ===================================================================================
        # EXPORT POWER POINT SLIDE
        # ===================================================================================
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        st.download_button(
            "⬇️ Download OCA Financial Report",
            data=pptx_io,
            file_name=f"OCA_Financial_Report_{datetime.now():%Y%m%d}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )











 # \\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\\ GENERATE PDF



# ================================================================ PART 2 PDF
    # ===================================================================================
    # 📄 EXECUTIVE-LEVEL OCA PDF REPORT (BEAUTIFIED)
    # ===================================================================================
    # ===================================================================================
    # 📄 FINAL PROFESSIONAL PDF (LANDSCAPE + CLEAN TABLES + LEGENDS)
    # ===================================================================================
    # 📄 FINAL PROFESSIONAL PDF (FULL TABLE + NO OVERFLOW)
    # ===================================================================================
    # 📄 FINAL PROFESSIONAL PDF (FULL TABLE + SAFE MONEY FORMAT)
    # ===================================================================================
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        Image, PageBreak
    )
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch

    import matplotlib.pyplot as plt
    import io
    import os
    import datetime
    import pandas as pd

    # ===================================================================================
    # 🔥 SAFE MONEY FORMATTER (FIXES YOUR ERROR)
    # ===================================================================================
    def safe_money(val):
        try:
            return f"${float(str(val).replace('$', '').replace(',', '')):,.2f}"
        except:
            return "$0.00"

    # ===================================================================================
    # 🔥 TEXT WRAP
    # ===================================================================================
    def wrap_text(text, max_len=28):
        words = str(text).split()
        lines, current = [], ""

        for w in words:
            if len(current + " " + w) <= max_len:
                current += " " + w
            else:
                lines.append(current.strip())
                current = w
        lines.append(current.strip())
        return "\n".join(lines)

    # ===================================================================================
    # 🔥 FULL TABLE BUILDER (ALL COLUMNS SAFE)
    # ===================================================================================
    def build_full_table(df):

        columns = [
            "COUPLE_ID",
            "NAME",
            "MONTHLY_TOTAL",
            "TOTAL_PAID",
            "EXPECTED",
            "BALANCE",
            "MONTHS_BEHIND",
            "STATUS"
        ]

        data = [columns]

        for _, r in df.iterrows():
            data.append([
                int(r["COUPLE_ID"]) if pd.notna(r["COUPLE_ID"]) else "",
                wrap_text(r["NAME"], 30),
                safe_money(r["MONTHLY_TOTAL"]),
                safe_money(r["TOTAL_PAID"]),
                safe_money(r["EXPECTED"]),
                safe_money(r["BALANCE"]),
                int(r["MONTHS_BEHIND"]) if pd.notna(r["MONTHS_BEHIND"]) else 0,
                r["STATUS"]
            ])

        col_widths = [
            0.7 * inch,
            3.3 * inch,
            1.1 * inch,
            1.1 * inch,
            1.1 * inch,
            1.1 * inch,
            1.0 * inch,
            1.0 * inch
        ]

        table = Table(data, colWidths=col_widths, repeatRows=1)

        table.setStyle(TableStyle([
            ("GRID", (0, 0), (-1, -1), 0.25, colors.black),

            ("BACKGROUND", (0, 0), (-1, 0), colors.darkblue),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),

            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),

            ("ALIGN", (0, 0), (0, -1), "CENTER"),
            ("ALIGN", (2, 1), (-2, -1), "RIGHT"),
            ("ALIGN", (-1, 1), (-1, -1), "CENTER")
        ]))

        return table

    # ===================================================================================
    # 📄 DOWNLOAD BUTTON
    # ===================================================================================
    st.subheader("📄 Download Executive Financial PDF")

    if st.button("📄 Download Executive PDF Report"):

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=landscape(letter))
        styles = getSampleStyleSheet()
        elements = []

        # ===================================================================================
        # TITLE
        # ===================================================================================
        elements.append(Paragraph("Owerri Cultural Association (OCA)", styles["Title"]))
        elements.append(Paragraph("Financial Report", styles["Heading1"]))
        elements.append(Paragraph(datetime.datetime.now().strftime("%B %Y"), styles["Heading2"]))
        elements.append(Spacer(1, 15))

        elements.append(Paragraph("Financial Secretary: Mr. Ugo Harris Opara", styles["Normal"]))
        elements.append(Paragraph("Assistant FS: Mrs. Pamela Uwakwe", styles["Normal"]))
        elements.append(Paragraph("Treasurer: Lady Ogechi Iwu", styles["Normal"]))

        elements.append(PageBreak())

        # ===================================================================================
        # 🟢 PAID IN FULL
        # ===================================================================================
        elements.append(Paragraph("🟢 Paid In Full", styles["Heading2"]))
        elements.append(build_full_table(paid_full))
        elements.append(PageBreak())

        # ===================================================================================
        # 🟡 CURRENT
        # ===================================================================================
        elements.append(Paragraph("🟡 Current (≤ 3 Months Behind)", styles["Heading2"]))
        elements.append(build_full_table(current))
        elements.append(PageBreak())

        # ===================================================================================
        # 🔴 BEHIND
        # ===================================================================================
        elements.append(Paragraph("🔴 Behind (3+ Months)", styles["Heading2"]))
        elements.append(build_full_table(behind))
        elements.append(PageBreak())

        # ===================================================================================
        # 📊 SUMMARY (SAFE FORMAT)
        # ===================================================================================
        elements.append(Paragraph("📊 Summary", styles["Heading2"]))

        summary_data = [["Months Behind", "Count", "Total Balance"]]

        for _, r in summary.iterrows():
            summary_data.append([
                int(r["MONTHS_BEHIND"]),
                int(r["COUNT"]),
                safe_money(r["TOTAL_BALANCE"])
            ])

        summary_table = Table(summary_data)
        summary_table.setStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.black)])

        elements.append(summary_table)
        elements.append(PageBreak())

        # ===================================================================================
        # 📈 CASHFLOW CHART
        # ===================================================================================
        elements.append(Paragraph("📈 Monthly Cashflow Analysis", styles["Heading2"]))

        fig, ax = plt.subplots(figsize=(10, 5))
        ax.plot(monthly_df["MONTH_STR"], monthly_df["Money In ($)"], marker='o', label="Money In")
        ax.plot(monthly_df["MONTH_STR"], monthly_df["Money Out ($)"], marker='o', label="Money Out")
        ax.legend()
        plt.xticks(rotation=45)

        img = io.BytesIO()
        fig.savefig(img, format="png", dpi=200)
        img.seek(0)

        elements.append(Image(img, width=700, height=350))
        elements.append(PageBreak())

        # ===================================================================================
        # 📆 ANNUAL BENCHMARK
        # ===================================================================================
        elements.append(Paragraph("📆 Annual Benchmark", styles["Heading2"]))

        fig, ax = plt.subplots(figsize=(14, 7))
        vals = [expected_full_year, total_received]

        bars = ax.bar(["Expected", "Received"], vals)

        for b in bars:
            h = b.get_height()
            ax.text(b.get_x() + b.get_width() / 2, h, f"${h:,.0f}", ha='center')

        img = io.BytesIO()
        fig.savefig(img, format="png", dpi=200)
        img.seek(0)

        elements.append(Image(img, width=750, height=400))

        # ===================================================================================
        # BUILD PDF
        # ===================================================================================
        doc.build(elements)
        buffer.seek(0)

        st.download_button(
            "⬇️ Download Final PDF",
            data=buffer,
            file_name=f"OCA_Report_{datetime.datetime.now():%Y%m%d}.pdf",
            mime="application/pdf"
        )













# ====================================================================================================
# ATTENDANCE ENDED HERE
# ======================================================================================================

        # ====================================================================================================
        # 🔥 DELETE MEMBER FROM DATABASE (ADMIN CONTROL)
        # ====================================================================================================

        # ======================================================================
        # 🔥 DELETE MEMBER SECTION (PASTE HERE — SAME INDENT AS ABOVE SECTIONS)
        # ======================================================================
        st.divider()
        st.subheader("❌ Delete Member From OCA Database")

        members_df = st.session_state.members_df

        if not members_df.empty:
            member_to_delete = st.selectbox(
                "Select a member to delete from the roster & attendance history:",
                options=members_df["name"].tolist(),
                key="delete_member_select"
            )

            if st.button("🗑️ Permanently Delete Member"):
                try:
                    import sqlite3
                    conn = sqlite3.connect("oca_attendance.db")
                    cur = conn.cursor()

                    cur.execute("""
                           DELETE FROM attendance
                           WHERE member_id = (SELECT id FROM members WHERE name = ?)
                       """, (member_to_delete,))

                    cur.execute("DELETE FROM members WHERE name = ?", (member_to_delete,))
                    conn.commit()
                    conn.close()

                    st.success(f"Member '{member_to_delete}' has been deleted.")

                    # Refresh UI
                    updated = get_members()
                    st.session_state.members_df = updated
                    st.session_state.roster = updated["name"].sort_values().tolist()
                    _refresh_attendance_from_db()

                except Exception as e:
                    st.error(f"Error deleting member: {e}")

        else:
            st.info("No members available to delete.")





# ===================================================================================
# MAIN APP
# # # ===================================================================================
# def main():
#     financial_ui()
#     oca_dues_ui()

# ======================================================================================================================================================
#                RUN APPLICATION
# ======================================================================================================================================================
if __name__ == "__main__":
    main()




















#
#
# # ---------------------------------------------------------------------
# # 2️⃣ CSV ANALYZER PAGE
# # ---------------------------------------------------------------------
# def csv_analyzer_page():
#     st.title("📊 CSV File Analyzer")
#     uploaded_file = st.file_uploader("Upload a CSV file", type=["csv"])
#
#     if uploaded_file is not None:
#         df = pd.read_csv(uploaded_file)
#         st.subheader("Dataset Preview")
#         st.write(df.head())
#
#         st.subheader("Basic Statistics")
#         st.write(df.describe())
#
#         st.subheader("Data Visualization")
#         numeric_columns = df.select_dtypes(include=['number']).columns.tolist()
#         if numeric_columns:
#             selected_column = st.selectbox("Select column to visualize", numeric_columns)
#             # Histogram
#             fig, ax = plt.subplots()
#             df[selected_column].hist(bins=20, edgecolor='black', ax=ax)
#             st.pyplot(fig)
#             # Boxplot
#             fig, ax = plt.subplots()
#             df[[selected_column]].boxplot(ax=ax)
#             st.pyplot(fig)
#         else:
#             st.warning("No numeric columns found for visualization.")
#
# # ---------------------------------------------------------------------
# # 3️⃣ ATTENDANCE PAGE
# # ---------------------------------------------------------------------
# def attendance_page():
#     st.title("🧾 OCA Attendance Tracker")
#     st.info("This section manages attendance, reporting, and visualization.")
#
#     # Example placeholder
#     st.write("➡️ Full attendance UI from your previous code will go here.")
#     st.caption("Tip: import your existing functions like `_init_state()` and `attendance_ui()` here.")
#
# # ---------------------------------------------------------------------
# # 4️⃣ MAIN APP MENU
# # ---------------------------------------------------------------------
# def main():
#     if not st.session_state.get('logged_in', False):
#         login_page()
#         return
#
#     # Sidebar Navigation
#     menu = st.sidebar.radio(
#         "📚 Main Menu",
#         ["🏠 Home", "📊 CSV Analyzer", "🧾 Attendance Tracker", "⚙️ Logout"]
#     )
#
#     if menu == "🏠 Home":
#         st.title(f"👋 Welcome, {st.session_state.username}!")
#         st.markdown("""
#         **Sybest LLC Internal Dashboard**
#         - Analyze CSV files
#         - Track Attendance
#         - Add new tools and modules anytime
#         """)
#     elif menu == "📊 CSV Analyzer":
#         csv_analyzer_page()
#     elif menu == "🧾 Attendance Tracker":
#         attendance_page()
#     elif menu == "⚙️ Logout":
#         st.session_state.logged_in = False
#         st.session_state.username = ''
#         st.experimental_rerun()
#
# # ---------------------------------------------------------------------
# if __name__ == "__main__":
#     main()


# =============================================== NEW STARTING HERE 10/31/2025 HOW


# XXXXXXXXXXXXXXXXXXXXXXXXXXXXXX now NOW NOW NOW














#
#
#
# # ]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]
# #
# # # CHATGPT: https://chatgpt.com/c/67d0dead-3d90-800f-ab0e-8dbb153d28cd
# # #
# import streamlit as st
# import pandas as pd
# import matplotlib.pyplot as plt
#
# # Streamlit App Title
# st.title("CSV File Analyzer")
#
# # File Uploader
# uploaded_file = st.file_uploader("Upload a CSV file", type=["csv"])
#
# if uploaded_file is not None:
#     # Read CSV file into DataFrame
#     df = pd.read_csv(uploaded_file)
#
#     # Display dataset
#     st.subheader("Dataset Preview")
#     st.write(df.head())
#
#     # Show basic statistics
#     st.subheader("Basic Statistics")
#     st.write(df.describe())
#
#     # Show column selection for visualization
#     st.subheader("Data Visualization")
#     numeric_columns = df.select_dtypes(include=['number']).columns.tolist()
#
#     if len(numeric_columns) > 0:
#         selected_column = st.selectbox("Select a column to visualize", numeric_columns)
#
#         # Histogram
#         st.subheader(f"Histogram of {selected_column}")
#         fig, ax = plt.subplots()
#         df[selected_column].hist(bins=20, edgecolor='black', ax=ax)
#         st.pyplot(fig)
#
#         # Box Plot
#         st.subheader(f"Boxplot of {selected_column}")
#         fig, ax = plt.subplots()
#         df[[selected_column]].boxplot(ax=ax)
#         st.pyplot(fig)
#     else:
#         st.write("No numeric columns available for visualization.")
#
# # ================================================================================== MAIN OCA RPASTER

#
# if __name__ == "__main__":
#     attendance_ui()


# # ]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]]] PART 444
#
# st.divider()
#
# # ---- Attendance Tracker (Hardcoded + Easy Editor + Jan–Nov Report + Roster Import/Export + Visuals) ----
# import streamlit as st
# import pandas as pd
# from datetime import date
# from io import BytesIO
# import altair as alt
#

# #
# STATUS_OPTIONS = ["Present", "Absent", "Excused Absent"]
# ANNUAL_MEETINGS = 10                  # Feb–Nov (Jan/Dec no meetings)
# MEETING_MONTHS = list(range(2, 12))   # 2..11 (Feb..Nov)
# MONTH_LABELS = ["FEB","MAR","APR","MAY","JUN","JUL","AUG","SEP","OCT","NOV"]
#
# # ---------------- Hardcoded roster ----------------
# DEFAULT_ROSTER = [
#     "Acholonu Kelechi & Nkiru (Mr. & Mrs.)",
#     "Akali Pet (Mrs.)",
#     "Amaechi Richard & Christina (Mr. & Mrs.)",
#     "Anyanwu Rosemary (Mrs)",
#     "Anyaso Hyacinth & Lydia (Chief & Lolo)",
#     "Chiagoro Chukwuma (Mr.)",
#     "Diala Emma & Joyce (Mr. & Mrs.)",
#     "Duru Gloria (Ms.)",
#     "Echibe Lucy (Ms.)",
#     "Egu Agnes (Mrs)",
#     "Egu Ernest (Mr. & Mrs.)",
#     "Egu Joy (Mrs)",
#     "Egu Mama Rose (Chief)",
#     "Egu Robert (Mr. & Mrs)",
#     "Eke Alexander (Mr. & Mrs.)",
#     "Eke Beatrice (Mrs.)",
#     "Eke Damian & Uloma (Chief & Lolo)",
#     "Ekeh George & Tina (Mr. & Mrs.)",
#     "Emeziem Adanma (Mrs.)",
#     "Ibe Prince & Chinwe (Mr. & Mrs.)",
#     "Ihejeto Anthony (Mr.)",
#     "Iwu Loretta (Ms.)",
#     "Iwuagwu Johnny (Mr. & Mrs.)",
#     "Mbara Tom (Mr. & Mrs.)",
#     "Njoku Stanley & Amaka (Mr. and Mrs)",
#     "Nwadibia Joannes Ubanwa (Ms.)",
#     "Nwadike Constance (Ms.)",
#     "Nwaogu Augustine (Mr. & Mrs.)",
#     "Nwaulu (Mr & Mrs)",
#     "Obichere Vivian (Mrs.)",
#     "Obilor Adolph (Chief & Dr. Lolo)",
#     "Obilor Innocent (Mr.)",
#     "Obilor Peter (Mr.)",
#     "Obinna Casmir (Mr.)",
#     "Ogbuehi Ethelbert (Mr.)",
#     "Ogwudire Kingsley & Ijeoma (Mr. & Mrs.)",
#     "Oha Linus & Liz (Mr. & Mrs.)",
#     "Oha Paul (Mr. & Mrs.)",
#     "Ohwobete Augustine and Grace (Dr. & Mrs.)",
#     "Ojibe Julie (Mrs)",
#     "Okoroafor Canon Progress (Rev)",
#     "Onyeagocha Rose (Lolo)",
#     "Onyejekwe Fidelia Chinyere (Mrs.)",
#     "Onyeneke Theophylus (Mr. & Mrs.)",
#     "Onyeukwu Damian (Chief)",
#     "Onyewuenyi Ethelbert R. (Mr.)",
#     "Opara Ugo Harris (Mr. & Mrs.)",
#     "Oparaocha Emeka & Joyce (Mr. & Mrs.)",
#     "Orji Christopher and Jacinta (Mr. & Mrs.)",
#     "Orji Ike F (Chief & Lolo)",
#     "Osuala Judy and Celine (Mr. & Mrs.)",
#     "Udeji Wilson",
#     "Ugorji Obi and Chinonye (Mr. & Mrs.)",
#     "Unaji Gibson & Victoria (Mr. & Mrs.)",
#     "Unanwa Christian (Dr. & Mrs)",
#     "Uwakwe Austin (Mr. & Mrs)",
#     "Uzohuo Uzoma David (Mr. & Mrs)",
#     "Uzoma Chinonye (Mr. & Mrs.)",
#     "Uzoma Christina (Ms.)",
#     "Yoko-Uzoma Okey (Mr. & Mrs.)",
# ]
#
# # ---------------- Session state ----------------
# def _init_state():
#     if "attendance_df" not in st.session_state:
#         st.session_state.attendance_df = pd.DataFrame(columns=["Date","Name","Status","Notes"])
#         # Make Notes explicitly text-friendly
#         st.session_state.attendance_df["Notes"] = st.session_state.attendance_df["Notes"].astype(object)
#     if "roster" not in st.session_state:
#         st.session_state.roster = sorted(pd.unique(DEFAULT_ROSTER).tolist())
#
# # ---------------- Instructions ----------------
# def _status_instructions():
#     st.info(
#         "**Allowed Attendance values**\n\n"
#         "• **Present**  (also: **P**, **Y**, **Yes**, **1**, **True**)\n"
#         "• **Absent**   (also: **A**, **N**, **No**, **0**, **False**)\n"
#         "• **Excused Absent**  (also: **E**, **EA**, **Excused**, **A.E.**)\n\n"
#         "_Case-insensitive. Dashes/underscores/dots OK — e.g. `excused-absent`, `excused_absent`, `A.E.`._"
#     )
#
# # ---------------- Helpers ----------------
# def _upsert_record(att_date, name, status, note=""):
#     """Insert or update a single record, always storing Notes as text."""
#     # normalize note -> text (avoid float/NaN)
#     if note is None or (isinstance(note, float) and pd.isna(note)):
#         note = ""
#     else:
#         note = str(note)
#
#     df = st.session_state.attendance_df
#     att_date = pd.to_datetime(att_date)
#     mask = (pd.to_datetime(df["Date"]) == att_date) & (df["Name"] == name)
#     if mask.any():
#         st.session_state.attendance_df.loc[mask, ["Status","Notes"]] = [status, note]
#     else:
#         st.session_state.attendance_df = pd.concat(
#             [df, pd.DataFrame([{"Date": att_date, "Name": name, "Status": status, "Notes": note}])],
#             ignore_index=True
#         )
#
# def _month_status(sub: pd.DataFrame) -> str:
#     """Summarize a person's month: Present > Excused Absent > Absent; blank if no records."""
#     if sub.empty: return ""
#     if (sub["Status"] == "Present").any(): return "Present"
#     if (sub["Status"] == "Excused Absent").any(): return "Excused Absent"
#     if (sub["Status"] == "Absent").any(): return "Absent"
#     return ""
#
# def _normalize_status(raw) -> str:
#     """Robust mapping to Present / Absent / Excused Absent (handles short forms & typos, incl. A.E.)."""
#     s = "" if pd.isna(raw) else str(raw).strip().lower()
#     s = s.replace(".", " ").replace("-", " ").replace("_", " ")
#     s_compact = s.replace(" ", "")
#     mapping = {
#         # Present
#         "present": "Present", "p": "Present", "pres": "Present",
#         "y": "Present", "yes": "Present", "1": "Present", "true": "Present",
#         # Absent
#         "absent": "Absent", "a": "Absent", "n": "Absent", "no": "Absent",
#         "0": "Absent", "false": "Absent", "nil": "Absent",
#         # Excused Absent
#         "excused absent": "Excused Absent", "excused": "Excused Absent",
#         "e": "Excused Absent", "ea": "Excused Absent", "a e": "Excused Absent",
#         "a.e.": "Excused Absent", "absent excused": "Excused Absent",
#         "excusedabsence": "Excused Absent", "ae": "Excused Absent",
#     }
#     if s in mapping: return mapping[s]
#     if s_compact in mapping: return mapping[s_compact]
#     # common misspellings
#     if "abesent" in s or "abesent" in s_compact: return "Absent"
#     if "excused abesent" in s or "excusedabesent" in s_compact: return "Excused Absent"
#     return ""
#
# def _find_col(cols, options):
#     """Find a column in cols (case-insensitive) matching any of options, or substring match."""
#     low = {c.lower(): c for c in cols}
#     for opt in options:
#         if opt.lower() in low:
#             return low[opt.lower()]
#     for c in cols:
#         lc = c.lower()
#         if any(opt.lower() in lc for opt in options):
#             return c
#     return None
#
# def _download_rollcall_csv(roster, att_date) -> bytes:
#     df = pd.DataFrame({
#         "Name": roster,
#         "Date": pd.to_datetime(att_date).date(),
#         "Attendance Status": ["" for _ in roster],
#         "Notes": ["" for _ in roster],
#     })
#     return df.to_csv(index=False).encode("utf-8")
#
# def _download_rollcall_xlsx(roster, att_date) -> bytes | None:
#     """Write a simple XLSX (no dropdowns). Returns bytes or None if engine missing."""
#     try:
#         out = BytesIO()
#         with pd.ExcelWriter(out, engine="openpyxl") as writer:
#             df = pd.DataFrame({
#                 "Name": roster,
#                 "Date": pd.to_datetime(att_date).date(),
#                 "Attendance Status": ["" for _ in roster],
#                 "Notes": ["" for _ in roster],
#             })
#             df.to_excel(writer, index=False, sheet_name="Roll Call")
#         return out.getvalue()
#     except Exception:
#         return None
#
# def _ingest_completed_roster(file, att_date, mark_unlisted_absent=False):
#     """Read CSV/XLSX with Name + Attendance Status (+ optional Notes) and apply to selected date."""
#     try:
#         if file.name.lower().endswith(".csv"):
#             df = pd.read_csv(file)
#         else:
#             df = pd.read_excel(file, sheet_name=0)
#     except Exception as e:
#         st.error(f"Could not read file: {e}")
#         return
#
#     name_col = _find_col(df.columns, ["Name"])
#     status_col = _find_col(df.columns, ["Attendance Status", "Status", "Attendance"])
#     notes_col = _find_col(df.columns, ["Notes", "Remark", "Comments"])
#
#     if not name_col or not status_col:
#         st.error("Upload must include columns for Name and Attendance Status (or Status/Attendance).")
#         return
#
#     # Normalize values
#     df[name_col] = df[name_col].astype(str).str.strip()
#     df["__status"] = df[status_col].map(_normalize_status)
#
#     bad = df["__status"] == ""
#     if bad.any():
#         st.error("Some rows have invalid status. Allowed: Present, Absent, Excused Absent (also P/A/E/Yes/No/A.E.).")
#         st.dataframe(df.loc[bad, [name_col, status_col]], use_container_width=True, hide_index=True)
#         return
#
#     names_in_file = set()
#     applied = 0
#     for _, row in df.iterrows():
#         nm = row[name_col].strip()
#         if not nm:
#             continue
#         names_in_file.add(nm)
#         status = row["__status"]
#         note_raw = row[notes_col] if notes_col else ""
#         note = "" if pd.isna(note_raw) else str(note_raw)
#         # If a new name appears, auto-add to roster
#         if nm not in st.session_state.roster:
#             st.session_state.roster.append(nm)
#             st.session_state.roster = sorted(pd.unique(st.session_state.roster).tolist())
#         _upsert_record(att_date, nm, status, note)
#         applied += 1
#
#     # Optionally mark everyone else (not in upload) as Absent
#     if mark_unlisted_absent:
#         for nm in st.session_state.roster:
#             if nm not in names_in_file:
#                 _upsert_record(att_date, nm, "Absent", "")
#
#     st.success(f"Applied {applied} updates for {pd.to_datetime(att_date).date()}.")
#
# # ---------------- Report helpers ----------------
# def _build_report_table(year: int) -> pd.DataFrame:
#     """
#     Jan..Nov matrix per Name with Totl Present and Percentage Present (denominator = 10).
#     Jan/Dec display if present in data, but are NOT counted in the denominator.
#     """
#     df = st.session_state.attendance_df.copy()
#     today_str = f"{date.today().month}/{date.today().day}/{date.today().year}"
#
#     if df.empty:
#         base = pd.DataFrame({"Name": st.session_state.roster})
#         for m in MONTH_LABELS:
#             base[m] = ""
#         base["Totl Present"] = 0
#         base["Percentage Present"] = "0%"
#         base.insert(1, "Month", today_str)
#         return base
#
#     df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
#     df = df[df["Date"].dt.year == year]
#
#     rows = []
#     for name in st.session_state.roster:
#         row = {"Name": name, "Month": today_str}
#         tot_present = 0
#         for idx, label in enumerate(MONTH_LABELS, start=1):  # 1..11 -> JAN..NOV
#             month_num = idx
#             sub = df[(df["Name"] == name) & (df["Date"].dt.month == month_num)]
#             status = _month_status(sub)
#             row[label] = status
#             if month_num in MEETING_MONTHS and status == "Present":
#                 tot_present += 1
#         row["Totl Present"] = tot_present
#         pct = round(tot_present / ANNUAL_MEETINGS * 100) if ANNUAL_MEETINGS else 0
#         row["Percentage Present"] = f"{pct}%"
#         rows.append(row)
#
#     report = pd.DataFrame(rows, columns=["Name","Month"] + MONTH_LABELS + ["Totl Present","Percentage Present"])
#     return report
#
# def _summary_cards(roster: list[str]):
#     def _is_couple(n: str) -> bool:
#         s = n.strip().lower()
#         return (" & " in s) or (" and " in s)
#     total_entries = len(roster)
#     couples = sum(1 for n in roster if _is_couple(n))
#     singles = total_entries - couples
#     total_individuals = couples*2 + singles
#     st.markdown(f"""
#         <style>
#             .cards {{display:flex; gap:16px; flex-wrap:wrap; margin: 4px 0 18px;}}
#             .card {{flex:1 1 240px; border:1px solid #e7e7e9; background:#fff; border-radius:16px;
#                     padding:16px 18px; box-shadow:0 2px 8px rgba(0,0,0,0.04);}}
#             .label {{font-size:13px; color:#666; margin-bottom:6px;}}
#             .value {{font-size:28px; font-weight:700; color:#111;}}
#         </style>
#         <div class="cards">
#           <div class="card"><div class="label">Total entries</div><div class="value">{total_entries}</div></div>
#           <div class="card"><div class="label">Singles</div><div class="value">{singles}</div></div>
#           <div class="card"><div class="label">Couples</div><div class="value">{couples}</div></div>
#           <div class="card"><div class="label">Total individual members</div><div class="value">{total_individuals}</div></div>
#         </div>
#     """, unsafe_allow_html=True)
#
# # ---------------- Visual helpers ----------------
# def _daily_totals_long(df: pd.DataFrame) -> pd.DataFrame:
#     if df.empty:
#         return pd.DataFrame(columns=["Date","Status","Count"])
#     daily = (df.groupby(["Date","Status"]).size()
#                .unstack(fill_value=0)
#                .reindex(columns=STATUS_OPTIONS, fill_value=0))
#     daily = daily.sort_index()
#     daily_reset = daily.reset_index()
#     daily_reset["Date"] = pd.to_datetime(daily_reset["Date"]).dt.date
#     # melt to long
#     long_df = daily_reset.melt(id_vars=["Date"], value_vars=STATUS_OPTIONS,
#                                var_name="Status", value_name="Count")
#     return long_df, daily_reset
#
# def _perperson_annual_counts(year: int) -> pd.DataFrame:
#     """Counts for Feb–Nov of the given year, plus Percent Attendance numeric."""
#     df_year = st.session_state.attendance_df.copy()
#     if df_year.empty:
#         counts = pd.DataFrame(columns=["Name"] + STATUS_OPTIONS)
#         counts["Percent Attendance"] = 0.0
#         counts["Percent_Attendance_Num"] = 0.0
#         return counts
#     df_year["Date"] = pd.to_datetime(df_year["Date"], errors="coerce")
#     df_year = df_year[(df_year["Date"].dt.year == year) & (df_year["Date"].dt.month.between(2,11))]
#     counts = (
#         df_year.groupby(["Name","Status"]).size()
#                .unstack(fill_value=0)
#                .reindex(columns=STATUS_OPTIONS, fill_value=0)
#                .reindex(st.session_state.roster, fill_value=0)
#                .reset_index()
#     )
#     counts["Percent_Attendance_Num"] = (counts["Present"] / ANNUAL_MEETINGS * 100).round(1)
#     counts["Percent Attendance"] = counts["Percent_Attendance_Num"].round(0).astype(int).astype(str) + "%"
#     return counts
#
# def _plot_daily_totals_chart(long_df: pd.DataFrame):
#     if long_df.empty:
#         st.caption("No daily data to plot yet.")
#         return
#     chart = (
#         alt.Chart(long_df)
#         .mark_bar()
#         .encode(
#             x=alt.X("Date:O", title="Date"),
#             y=alt.Y("sum(Count):Q", title="Headcount"),
#             color=alt.Color("Status:N", sort=STATUS_OPTIONS),
#             tooltip=["Date","Status","Count"]
#         )
#         .properties(height=320)
#     )
#     st.altair_chart(chart, use_container_width=True)
#
# def _plot_perperson_chart(df_counts: pd.DataFrame, sort_by: str, top_n: int):
#     if df_counts.empty:
#         st.caption("No annual data to plot yet.")
#         return
#     if sort_by == "Present (count)":
#         dfp = df_counts.sort_values("Present", ascending=False).head(top_n)
#         x_field = "Present"
#         x_title = "Present (Feb–Nov)"
#     else:
#         dfp = df_counts.sort_values("Percent_Attendance_Num", ascending=False).head(top_n)
#         x_field = "Percent_Attendance_Num"
#         x_title = "Percent Attendance (%)"
#     chart = (
#         alt.Chart(dfp)
#         .mark_bar()
#         .encode(
#             x=alt.X(f"{x_field}:Q", title=x_title, scale=alt.Scale(domain=[0, 100]) if sort_by!="Present (count)" else alt.Undefined),
#             y=alt.Y("Name:N", sort="-x", title="Member"),
#             tooltip=["Name","Present","Absent","Excused Absent","Percent Attendance"]
#         )
#         .properties(height=24*len(dfp) + 60)
#     )
#     st.altair_chart(chart, use_container_width=True)
#
# # ---------------- UI ----------------
# def attendance_ui():
#     _init_state()
#     st.title("Attendance — Easy Taker + Jan–Nov Report + Roster Import/Export + Visuals")
#
#     # Summary
#     _summary_cards(st.session_state.roster)
#
#
#
#     # 1) Roster
#     st.subheader("1) Roster & Add New Member")
#     c1, c2 = st.columns([2,1])
#     with c1:
#         new_name = st.text_input("Add a new member (full name)")
#         if st.button("➕ Add member"):
#             nm = new_name.strip()
#             if nm:
#                 if nm not in st.session_state.roster:
#                     st.session_state.roster.append(nm)
#                     st.session_state.roster = sorted(st.session_state.roster)
#                     st.success(f"Added: {nm}")
#                 else:
#                     st.info("That name already exists.")
#             else:
#                 st.warning("Enter a name.")
#     with c2:
#         roster_csv = pd.DataFrame({"Name": st.session_state.roster}).to_csv(index=False)
#         st.download_button("Download roster.csv", roster_csv, "roster.csv", "text/csv")
#
#     st.divider()
#     st.subheader('COMPREHENSIVE OCA MEMBERS LIST')
#     with st.expander('LIST OF OCA MEMBERS'):
#
#         st.dataframe(pd.DataFrame({"Name": st.session_state.roster}),
#                      use_container_width=True, hide_index=True)
#
#     st.divider()
#     st.subheader('TAKE MEETING ATTENDANCE')
#
#     # 2) Take Attendance (beautified editor)
#     st.subheader("2) Take Attendance")
#     _status_instructions()
#
#     if not st.session_state.roster:
#         st.info("Add members to your roster first.")
#         return
#
#     # Selected date defaults to TODAY
#     cc1, cc2, cc3, cc4 = st.columns([2,2,2,1])
#     with cc1:
#         att_date = st.date_input("Attendance date", value=date.today())
#     with cc2:
#         default_status = st.selectbox("Default for 'Mark All'", STATUS_OPTIONS, index=0)
#     with cc3:
#         mark_all = st.button("Mark All with default")
#     with cc4:
#         clear_btn = st.button("🧽 Clear this date")
#
#     if mark_all:
#         for nm in st.session_state.roster:
#             _upsert_record(att_date, nm, default_status, "")
#         st.success("Marked all with default.")
#
#     if clear_btn:
#         df = st.session_state.attendance_df
#         st.session_state.attendance_df = df[pd.to_datetime(df["Date"]) != pd.to_datetime(att_date)]
#         st.info("Cleared all entries for this date.")
#
#     # Editable grid for the selected date
#     df_today = st.session_state.attendance_df.copy()
#     df_today["Date"] = pd.to_datetime(df_today["Date"])
#     df_today = df_today[df_today["Date"].dt.date == pd.to_datetime(att_date).date()]
#
#     current = {row["Name"]: (row["Status"], row.get("Notes","")) for _, row in df_today.iterrows()}
#
#     # Ensure Notes are TEXT (no NaN -> float)
#     notes_clean = []
#     for n in st.session_state.roster:
#         v = current.get(n, (default_status, ""))[1]
#         v = "" if v is None or (isinstance(v, float) and pd.isna(v)) else str(v)
#         notes_clean.append(v)
#
#     edit_df = pd.DataFrame({
#         "Name":   st.session_state.roster,
#         "Status": [current.get(n, (default_status, ""))[0] for n in st.session_state.roster],
#         "Notes":  notes_clean,
#     })
#     edit_df["Notes"] = edit_df["Notes"].astype(str)  # enforce text dtype
#
#     edited = st.data_editor(
#         edit_df,
#         use_container_width=True,
#         hide_index=True,
#         num_rows="fixed",
#         column_config={
#             "Name": st.column_config.TextColumn(disabled=True),
#             "Status": st.column_config.SelectboxColumn(options=STATUS_OPTIONS),
#             "Notes": st.column_config.TextColumn(),
#         },
#         key=f"editor_{att_date}",
#     )
#     if st.button("💾 Save updates for this date"):
#         for _, r in edited.iterrows():
#             _upsert_record(att_date, r["Name"], r["Status"], r["Notes"])
#         st.success("Saved.")
#
#     st.divider()
#     # --- Download / Upload roster for physical roll-call ---
#     st.markdown("#### Roll-call file (download for today / upload completed)")
#     _status_instructions()
#
#     dcol1, dcol2, ucol = st.columns([1,1,2])
#     with dcol1:
#         st.download_button(
#             "Download roll-call CSV (today)",
#             data=_download_rollcall_csv(st.session_state.roster, att_date),
#             file_name=f"rollcall_{pd.to_datetime(att_date).date()}.csv",
#             mime="text/csv"
#         )
#     with dcol2:
#         xlsx_bytes = _download_rollcall_xlsx(st.session_state.roster, att_date)
#         if xlsx_bytes:
#             st.download_button(
#                 "Download roll-call Excel (today)",
#                 data=xlsx_bytes,
#                 file_name=f"rollcall_{pd.to_datetime(att_date).date()}.xlsx",
#                 mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
#             )
#         else:
#             st.caption("Excel export requires `openpyxl`. CSV works without extras.")
#
#     with ucol:
#         mark_unlisted_absent = st.checkbox("Mark names not listed in upload as Absent (for this date)", value=False)
#         uploaded = st.file_uploader(
#             "Upload completed roster (CSV or XLSX) with Name + Attendance Status (+ Notes optional)",
#             type=["csv","xlsx"]
#         )
#         if uploaded is not None:
#             _ingest_completed_roster(uploaded, att_date, mark_unlisted_absent)
#
#     st.divider()
#
#     # 3) Summary & Export
#     st.subheader("3) Summary & Export")
#
#     # Raw log (optional)
#     df = st.session_state.attendance_df.copy()
#     if df.empty:
#         st.info("No attendance recorded yet.")
#         return
#
#     with st.expander("Raw attendance log (filterable)"):
#         dff = df.sort_values(["Date","Name"]).copy()
#         dff["Date"] = pd.to_datetime(dff["Date"]).dt.date
#         # quick filters
#         fcol1, fcol2 = st.columns([2,2])
#         with fcol1:
#             raw_name_search = st.text_input("Search name (raw log)", "")
#         with fcol2:
#             raw_status_filter = st.multiselect("Filter status (raw log)", STATUS_OPTIONS, default=[])
#         dff2 = dff.copy()
#         if raw_name_search.strip():
#             dff2 = dff2[dff2["Name"].str.contains(raw_name_search.strip(), case=False, na=False)]
#         if raw_status_filter:
#             dff2 = dff2[dff2["Status"].isin(raw_status_filter)]
#         st.dataframe(dff2, use_container_width=True, hide_index=True)
#
#     # Annual report in your layout (Month uses TODAY)
#     st.write("### Annual Report (Jan–Nov layout, Percent = Present ÷ 10 × 100)")
#     year = pd.Timestamp.today().year
#     report = _build_report_table(year)
#     st.dataframe(report, use_container_width=True, hide_index=True)
#     st.download_button("Download report (CSV)", report.to_csv(index=False), "attendance_report.csv", "text/csv")
#
#     st.divider()
#     st.subheader('INDIVIDUAL ATTENDANCE SUMMARY EXPAND TO VIEW')
#
#     with st.expander('INDIVIDUAL ATTENDANCE SUMMARY - HIGHLY -LEVEL'):
#         st.divider()
#         # Per-person totals (Annual, fixed 10 meetings) + filters + chart
#         st.write("### Per-person totals (Annual, fixed 10 meetings) — with search, status filter & chart")
#         counts = _perperson_annual_counts(year)
#
#         # Filters for this section
#         f1, f2, f3 = st.columns([2,2,1])
#         with f1:
#             name_query = st.text_input("Search name", "")
#         with f2:
#             status_show = st.multiselect("Show members who have ANY of these statuses", STATUS_OPTIONS, default=[])
#         with f3:
#             top_n = st.slider("Top N (chart)", min_value=5, max_value=50, value=10, step=1)
#
#         counts_view = counts.copy()
#         if name_query.strip():
#             counts_view = counts_view[counts_view["Name"].str.contains(name_query.strip(), case=False, na=False)]
#         if status_show:
#             mask = counts_view[status_show].sum(axis=1) > 0
#             counts_view = counts_view[mask]
#
#         st.dataframe(counts_view[["Name","Present","Absent","Excused Absent","Percent Attendance"]],
#                      use_container_width=True, hide_index=True)
#
#         # Chart controls
#         sort_by = st.radio("Sort chart by", ["Present (count)", "Percent (%)"], horizontal=True)
#         sort_key = "Present (count)" if sort_by.startswith("Present") else "Percent Attendance"
#         _plot_perperson_chart(counts_view, "Present (count)" if sort_key=="Present (count)" else "Percent (%)", top_n)
#
#     st.divider()
#
#     # Daily totals (explicit Date column) + chart
#     st.write("### Daily totals (table + chart)")
#     long_df, daily_reset = _daily_totals_long(st.session_state.attendance_df)
#     st.dataframe(daily_reset, use_container_width=True, hide_index=True)
#     _plot_daily_totals_chart(long_df)
#
# if __name__ == "__main__":
#     attendance_ui()
#







































#
# # =========================================================================
# #   LOGIN OCA
# # ===========================================================================
#
# import streamlit as st
# import os
# from dotenv import load_dotenv
#
# # Load environment variables from .env file
# load_dotenv()
#
# # Retrieve user credentials from environment variables
# USER_CREDENTIALS = {
#     os.getenv('Stan3000'): os.getenv('Monday30@'),
#     os.getenv('Stan3000'): os.getenv('Monday30@')
# }
#
#
#
# def login():
#     st.title("SYBEST Login Page")
#
#     # Initialize session state variables
#     if 'logged_in' not in st.session_state:
#         st.session_state.logged_in = False
#     if 'username' not in st.session_state:
#         st.session_state.username = ''
#
#     # Create a form for login
#     with st.form(key='login_form'):
#         username = st.text_input("Username")
#         password = st.text_input("Password", type="password")
#         submit_button = st.form_submit_button(label="Login")
#
#     if submit_button:
#         if username in USER_CREDENTIALS and USER_CREDENTIALS[username] == password:
#             st.session_state.logged_in = True
#             st.session_state.username = username
#             st.success("Login successful!")
#         else:
#             st.error("Invalid username or password")
#
# # Call the login function to display the login page
# login()
#
# # Display content based on login status
# if st.session_state.logged_in:
#     st.write(f"Welcome, {st.session_state.username}!")
#     # Place additional authenticated user content here
# else:
#     st.write("Please log in to continue.")
# # #
# #
# #
# #
#
# #
# #
# #
# #
# #
# # # ========================================================================================================================================================
# # #                                                    PAGES
# # # =========================================================
# #
# # def main():
# #     st.sidebar.title('Data Selection')
# #     page_selection = st.sidebar.radio("Choose a dataset:", ["Home", "OCA CONSTITUTION",'EXAMINETICS INVOICING']) #'UNIFORM PRICING PROJECT'])
# #
# #     if page_selection == "Home":
# #         st.title("Welcome to the GovCon Dashboard SAM API")
# #         st.write("Select a dataset from the sidebar to begin analysis.")
# #
# #     # elif page_selection == "Sam Data":
# #     #     st.title("Analysis for Sam Data")
# #     #     st.write(r"DATA LOCATION: C:\Users\stans\OneDrive\Desktop\SYBEST LLC\DATASETS\SAM AND CAL-PROCEDURE_DATA")
# #     #     uploaded_file = st.file_uploader("Upload Sam Data Excel file", type=['xlsx'])
# #     #     if uploaded_file:
# #     #         df = pd.read_excel(uploaded_file, engine='openpyxl')
# #     #         sam_data_analysis(df)
# #     elif page_selection == "Sam Data":
# #         st.title("Analysis for Sam Data")
# #         st.write(r"DATA LOCATION: C:\Users\stans\OneDrive\Desktop\SYBEST LLC\DATASETS\INDIVIDUAL PROJECT DATASET")
# #         uploaded_file = st.file_uploader("Upload Sam Data Excel file", type=['csv'])
# #         if uploaded_file:
# #             import pandas as pd
# #             df = pd.read_csv(uploaded_file, encoding='ISO-8859-1')
# #             sam_data_analysis(df)
# #
# #
# #     elif page_selection == "OCA CONSTITUTION":
# #         st.title("Analysis for Sam Data")
# #         st.write(r"DATA LOCATION: C:\Users\stans\OneDrive\Desktop\SYBEST LLC\DATASETS\INDIVIDUAL PROJECT DATASET")
# #         uploaded_file = st.file_uploader("Upload Sam Data Excel file", type=['csv'])
# #         if uploaded_file:
# #             import pandas as pd
# #             df_con = pd.read_csv(uploaded_file, encoding='ISO-8859-1')
# #             sam_data_analysis(df_con)
#
# # ======================================================================================== PART 3
#
#
# import streamlit as st
# import pandas as pd
# from PyPDF2 import PdfReader
#

#
# # ======================== PAGE FUNCTIONS =========================
#
# def home_page():
#     st.title("🏠 Welcome to the GovCon Dashboard")
#     st.write("Select a page from the sidebar to begin analysis.")
#
#
# # ====================================================================================================
# # OCA CONSTITUTION
# # ======================================================================================================
#
# # def oca_constitution_page():
# #     st.title("📜 OCA Constitution Upload")
# #     uploaded_file = st.file_uploader("Upload OCA Constitution CSV file", type=['csv'])
# #     if uploaded_file:
# #         df = pd.read_csv(uploaded_file, encoding='ISO-8859-1')
# #         st.write("Preview of the Constitution Data:")
# #         st.dataframe(df)
# #
# # OCA CON PATH: C:\Users\stans\OneDrive\Desktop\OCA\OCA CONSITITUTION\OCA_constitution-Edited_02212025.docx
# #
# #
# # ==========================================================================================================PART 2
#
# from docx import Document
# import streamlit as st
# import os
# #
# # def load_docx_text_chunks(path, chunk_size=300):
# #     if not os.path.exists(path):
# #         st.error(f"File not found at: {path}")
# #         return []
# #
# #     doc = Document(path)
# #     full_text = ""
# #
# #     for para in doc.paragraphs:
# #         text = para.text.strip()
# #         if text:
# #             full_text += text + "\n"
# #
# #     # Break into pseudo-pages
# #     chunks = [full_text[i:i + chunk_size] for i in range(0, len(full_text), chunk_size)]
# #     return chunks
# #
# # def oca_constitution_page():
# #     st.title("📜 OCA Constitution Viewer and Keyword Search")
# #
# #     # Load file from static path
# #     file_path = r"C:\Users\stans\OneDrive\Desktop\OCA\OCA CONSITITUTION\OCA_constitution-Edited_02212025.docx"
# #     pages = load_docx_text_chunks(file_path)
# #
# #     if not pages:
# #         return
# #
# #     # Keyword Search
# #     keyword = st.text_input("🔍 Enter a keyword to search in the document")
# #     if keyword:
# #         matches = []
# #         for idx, page in enumerate(pages):
# #             if keyword.lower() in page.lower():
# #                 matches.append((idx + 1, page))
# #
# #         if matches:
# #             st.success(f"Found {len(matches)} match(es) for '{keyword}':")
# #             for page_num, content in matches:
# #                 st.markdown(f"### 📄 Page {page_num}")
# #                 highlighted = content.replace(keyword, f"**:red[{keyword}]**")
# #                 st.write(highlighted)
# #                 st.markdown("---")
# #         else:
# #             st.warning("No matches found.")
# #
# #     st.markdown("## 📖 Browse the Document Page by Page")
# #     selected_page = st.slider("Select Page Number", min_value=1, max_value=len(pages), value=1)
# #     st.markdown(f"### 📄 Page {selected_page}")
# #     st.text(pages[selected_page - 1])
#
#
# # ========================================================================
#
# # def load_docx_text_chunks(path, chunk_size=300):
# #     if not os.path.exists(path):
# #         st.error(f"File not found at: {path}")
# #         return []
# #
# #     doc = Document(path)
# #     full_text = ""
# #
# #     for para in doc.paragraphs:
# #         text = para.text.strip()
# #         if text:
# #             full_text += text + "\n"
# #
# #     # Break into text chunks
# #     chunks = [full_text[i:i + chunk_size] for i in range(0, len(full_text), chunk_size)]
# #     return chunks
# #
# # def oca_constitution_page():
# #     st.title("📜 OCA Constitution Viewer and Keyword Search")
# #
# #     file_path = r"C:\Users\stans\OneDrive\Desktop\OCA\OCA CONSITITUTION\OCA_constitution-Edited_02212025.pdf"
# #     pages = load_docx_text_chunks(file_path)
# #
# #     if not pages:
# #         return
# #
# #     keyword = st.text_input("🔍 Enter a keyword to search in the document")
# #
# #     default_page = 1
# #     match_highlight = None
# #
# #     if keyword:
# #         matches = []
# #         for idx, page in enumerate(pages):
# #             if keyword.lower() in page.lower():
# #                 matches.append((idx + 1, page))
# #
# #         if matches:
# #             st.success(f"Found {len(matches)} match(es) for '{keyword}':")
# #             for page_num, content in matches:
# #                 st.markdown(f"### 📄 Page {page_num}")
# #                 highlighted = content.replace(keyword, f"**:red[{keyword}]**")
# #                 st.write(highlighted)
# #                 st.markdown("---")
# #
# #             # Set first match as default to view
# #             default_page = matches[0][0]
# #             match_highlight = matches[0][1].replace(keyword, f"**:red[{keyword}]**")
# #         else:
# #             st.warning("No matches found.")
# #
# #     st.markdown("## 📖 Browse the Document Page by Page")
# #
# #     selected_page = st.slider("Select Page Number", min_value=1, max_value=len(pages), value=default_page)
# #     st.markdown(f"### 📄 Page {selected_page}")
# #
# #     # If current page is the matched page, show highlighted, else plain
# #     if match_highlight and selected_page == default_page:
# #         st.markdown(match_highlight)
# #     else:
# #         st.text(pages[selected_page - 1])
#
# # To run inside Streamlit
# # if __name__ == "__main__":
# # =================================================================== PART 2
#
# import os
# import fitz  # PyMuPDF
# import streamlit as st
# import re
#
# def load_pdf_text_chunks(path):
#     if not os.path.exists(path):
#         st.error(f"File not found at: {path}")
#         return []
#     doc = fitz.open(path)
#     pages = [page.get_text("text") for page in doc]
#     return pages
#
# def highlight_keyword(text, keyword):
#     # Escape keyword to avoid regex errors
#     escaped_keyword = re.escape(keyword)
#     pattern = re.compile(f"({escaped_keyword})", re.IGNORECASE)
#     # Replace with red-colored HTML span
#     highlighted_text = pattern.sub(r"<span style='color:red; font-weight:bold;'>\1</span>", text)
#     return highlighted_text
#
# def oca_constitution_page():
#     st.title("📜 OCA Constitution Viewer and Keyword Search")
#
#     file_path = r"C:\Users\stans\OneDrive\Desktop\OCA\OCA CONSITITUTION\OCA_constitution-Edited_02212025.pdf"
#     pages = load_pdf_text_chunks(file_path)
#
#     if not pages:
#         return
#
#     keyword = st.text_input("🔍 Enter a keyword to search in the document")
#     matches = []
#
#     if keyword:
#         for idx, content in enumerate(pages):
#             if keyword.lower() in content.lower():
#                 matches.append((idx + 1, content))  # Store 1-based page number
#
#         if matches:
#             st.success(f"Found {len(matches)} match(es) for '{keyword}':")
#             page_numbers = [page_num for page_num, _ in matches]
#             selected_page = st.selectbox("📄 Select a matched page to view", page_numbers)
#
#             selected_content = matches[page_numbers.index(selected_page)][1]
#             highlighted_html = highlight_keyword(selected_content, keyword)
#
#             st.markdown(f"### 📄 Page {selected_page}")
#             st.markdown(highlighted_html, unsafe_allow_html=True)
#         else:
#             st.warning("No matches found.")
#
#     else:
#         st.markdown("## 📖 Browse Full Document (No Keyword Entered)")
#         selected_page = st.slider("Select Page Number", 1, len(pages), 1)
#         st.markdown(f"### 📄 Page {selected_page}")
#         st.text(pages[selected_page - 1])
#
#
#
#
#
#
#
#
#
#
#
#
#
#
#
#
#
#
#
#
# ====================================================================================================================
#                                   OCA MEMBER START HERE
# =====================================================================================================================
# import streamlit as st
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt

def analyze_membership_debt():
    st.title("📊 OCA Membership Dues Analysis")

    uploaded_file = st.file_uploader("Upload Membership Excel File", type=["xlsx", "xls"])


    if uploaded_file:
        # Load Excel and skip first 3 rows
        df = pd.read_excel(uploaded_file, skiprows=4)

        # Display column headers for debug
        st.markdown("### 🧩 Detected Columns:")
        st.write(df.columns.tolist())

        # Normalize columns
        df.columns = df.columns.str.strip().str.lower()

        # Auto-detect important columns
        date_col = next((col for col in df.columns if 'date' in col), None)
        name_col = next((col for col in df.columns if 'name' in col), None)
        bal_col = next((col for col in df.columns if 'open' in col and 'bal' in col), None)

        if not all([date_col, name_col, bal_col]):
            st.error("❌ Missing required columns: 'date', 'name', or 'open balance'.")
            return

        # Rename for easier handling
        df = df.rename(columns={date_col: "Date", name_col: "Name", bal_col: "Open Balance"})

        # Clean and prepare data
        df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
        df = df.dropna(subset=["Date", "Name"])
        df["Year"] = df["Date"].dt.year
        df["Open Balance"] = pd.to_numeric(df["Open Balance"], errors="coerce").fillna(0)

        # Filter for years 2021–2023
        df_filtered = df[df["Year"].isin([2021, 2022, 2023, 2024])]

        # Build list of owing members per year
        owing_by_year = {
            year: df_filtered[(df_filtered["Year"] == year) & (df_filtered["Open Balance"] > 0)]["Name"].unique().tolist()
            for year in [2021, 2022, 2023, 2024]
        }

        st.divider()

# ============================================================
#  Members Who Owed in 2021, 2022, and 2023
# =============================================================

        # Identify members owing in all 3 years
        owing_all_years = list(set(owing_by_year[2021]) & set(owing_by_year[2022]) & set(owing_by_year[2023]) & set(owing_by_year[2024]))

        st.subheader("👥 Members Who Owed in 2021, 2022, and 2023")
        if owing_all_years:
            st.write(pd.DataFrame(owing_all_years, columns=["Name"]))
        else:
            st.success("✅ No members owed across all three years.")

        # Total unique members by year
        unique_members = df_filtered.groupby("Year")["Name"].nunique().to_dict()

        st.divider()

        # 📅 Unique Members and Owing Counts by Year with Total Amount Owed
        st.subheader("📅 Unique Members, Owing Counts & Totals by Year")

        col1, col2, col3, col4 = st.columns(4)
        colors = ["#e1f0ff", "#e1f0ff", "#e1f0ff"]  # Light blue background

        for i, (year, col) in enumerate(zip([2021, 2022, 2023, 2024], [col1, col2, col3, col4])):
            unique_count = unique_members.get(year, 0)
            owing_names = owing_by_year[year]

            # Calculate total owed that year
            total_owed = df_filtered[
                (df_filtered["Year"] == year) &
                (df_filtered["Name"].isin(owing_names))
                ]["Open Balance"].sum()

            col.markdown(f"""
                <div style="background-color: {colors[i]}; padding: 20px; border-radius: 10px; box-shadow: 0 1px 4px rgba(0,0,0,0.1);">
                    <h4 style="color:#0a3d62;">📅 Year: {year}</h4>
                    <p style="font-size: 16px; color:#333;"><strong>Unique Members:</strong> {unique_count}</p>
                    <p style="font-size: 16px; color:#333;"><strong>Owing Members:</strong> {len(owing_names)}</p>
                    <p style="font-size: 16px; color:#333;"><strong>Total Owed:</strong> ${total_owed:,.2f}</p>
                </div>
            """, unsafe_allow_html=True)


        st.divider()
#
# # =============================================================
#      #   Highest Owing Member Each Year
# # ===============================================================
#
#
#         # 🧾 Highest Owing Member Each Year
#         st.subheader("💰 Highest Owing Member Each Year")
#
#         highest_owers = []
#
#         for year in [2021, 2022, 2023, 2024]:
#             sub_df = df_filtered[(df_filtered["Year"] == year) & (df_filtered["Open Balance"] > 0)]
#             if not sub_df.empty:
#                 top_row = sub_df.loc[sub_df["Open Balance"].idxmax()]
#                 highest_owers.append({
#                     "Year": year,
#                     "Name": top_row["Name"],
#                     "Amount": top_row["Open Balance"]
#                 })
#             else:
#                 highest_owers.append({
#                     "Year": year,
#                     "Name": "—",
#                     "Amount": 0.00
#                 })
#
#         # Display in nice table format
#         df_top_owers = pd.DataFrame(highest_owers)
#         df_top_owers["Amount"] = df_top_owers["Amount"].apply(lambda x: f"${x:,.2f}")
#         st.table(df_top_owers.rename(columns={
#             "Year": "📅 Year",
#             "Name": "👤 Name",
#             "Amount": "💸 Amount Owed"
#         }))
#
#         st.divider()
#
# # ========================================================================
#   # TOP5 5 PEOPLE
# # ===========================================================================
#         with st.expander('💰 Top 5 Highest Owing Members Each Year'):
#             # 💰 Top 5 Highest Owing Members Per Year
#             st.subheader("💰 Top 5 Highest Owing Members Each Year")
#
#             for year in [2021, 2022, 2023, 2024,]:
#                 st.markdown(f"### 📅 Year: {year}")
#
#                 sub_df = df_filtered[
#                     (df_filtered["Year"] == year) &
#                     (df_filtered["Open Balance"] > 0)
#                     ]
#
#                 if not sub_df.empty:
#                     top5 = sub_df[["Name", "Open Balance"]].groupby("Name").sum()
#                     top5 = top5.sort_values(by="Open Balance", ascending=False).head(5).reset_index()
#                     top5["Open Balance"] = top5["Open Balance"].apply(lambda x: f"${x:,.2f}")
#                     st.table(top5.rename(columns={"Name": "👤 Name", "Open Balance": "💸 Amount Owed"}))
#                 else:
#                     st.info("No members with outstanding balances for this year.")
#
#         # ============================================================
# # Owing Members and Amounts by Yea
# # =============================================================
#
#
#         # 📋 Table of owing members with amounts
#         st.subheader("📋 Owing Members and Amounts by Year")
#
#         def get_name_and_amount(df, year):
#             sub_df = df[(df["Year"] == year) & (df["Open Balance"] > 0)]
#             return [f"{row['Name']} (${row['Open Balance']:,.2f})" for _, row in sub_df.iterrows()]
#
#         data_2021 = get_name_and_amount(df_filtered, 2021)
#         data_2022 = get_name_and_amount(df_filtered, 2022)
#         data_2023 = get_name_and_amount(df_filtered, 2023)
#         data_2024 = get_name_and_amount(df_filtered, 2024)
#         max_rows = max(len(data_2021), len(data_2022), len(data_2023), len(data_2024))
#
#         df_owing_table = pd.DataFrame({
#             "2021": data_2021 + [""] * (max_rows - len(data_2021)),
#             "2022": data_2022 + [""] * (max_rows - len(data_2022)),
#             "2023": data_2023 + [""] * (max_rows - len(data_2023)),
#             "2024": data_2024 + [""] * (max_rows - len(data_2024)),
#         })
#
#         st.dataframe(df_owing_table)
#
# # ============================================================
#         # Members Owing Per Year
# #=============================================================
#         st.divider()
#         # 📈 Bar chart with annotations
#         st.subheader("📈 Members Owing Per Year")
#         fig, ax = plt.subplots()
#         years = [2021, 2022, 2023, 2024]
#         counts = [len(owing_by_year[year]) for year in years]
#
#         bars = ax.bar(years, counts, color='skyblue')
#         for bar in bars:
#             height = bar.get_height()
#             ax.annotate(f'{int(height)}', xy=(bar.get_x() + bar.get_width() / 2, height),
#                         xytext=(0, 3), textcoords="offset points",
#                         ha='center', va='bottom', fontsize=10)
#
#         ax.set_ylabel("Members Owing")
#         ax.set_xlabel("Year")
#         ax.set_title("Outstanding Dues by Year")
#         st.pyplot(fig)
#
#
# # ==================================================
# # OCA Members - Hardcoded, Editable, Deletable
# # ===========================================
#         st.subheader("📋 OCA Members – Manage Manually")
#
#         # Hardcoded list (only loads once)
#         default_members = [
#             {"Name": "John Doe", "Open Balance": 150.0},
#             {"Name": "Angela Nwosu", "Open Balance": 200.0},
#             {"Name": "Peter Okafor", "Open Balance": 100.0},
#             {"Name": "Chinwe Opara", "Open Balance": 80.0},
#             {"Name": "Samuel Obi", "Open Balance": 50.0},
#         ]
#
#         # Store in session state
#         if "oca_manual_members" not in st.session_state:
#             st.session_state.oca_manual_members = default_members.copy()
#
#         members = st.session_state.oca_manual_members
#
#         # Display editable list
#         for i, member in enumerate(members):
#             with st.expander(f"👤 {member['Name']} – ${member['Open Balance']:,.2f}"):
#                 col1, col2, col3 = st.columns([4, 2, 1])
#
#                 with col1:
#                     new_name = st.text_input(f"Edit Name #{i}", value=member["Name"], key=f"edit_name_{i}")
#                 with col2:
#                     new_amount = st.number_input(f"Edit Amount #{i}", value=member["Open Balance"], step=10.0, key=f"edit_amt_{i}")
#                 with col3:
#                     delete = st.button("🗑️ Delete", key=f"delete_{i}")
#
#                 if delete:
#                     st.session_state.oca_manual_members.pop(i)
#                     st.experimental_rerun()
#
#                 # Save changes
#                 member["Name"] = new_name
#                 member["Open Balance"] = new_amount
#
#         # Display table of current state
#         st.markdown("### ✅ Updated OCA Members List")
#         df_updated = pd.DataFrame(st.session_state.oca_manual_members)
#         df_updated["Open Balance"] = df_updated["Open Balance"].apply(lambda x: f"${x:,.2f}")
#         st.dataframe(df_updated)


# ================================================================================= PART 2
#
#
# import streamlit as st
# import pandas as pd
# import matplotlib.pyplot as plt
#
# def analyze_membership_debt():
#     st.title("📊 OCA Membership Dues Analysis")
#
#     uploaded_file = st.file_uploader("Upload Membership Excel File", type=["xlsx", "xls"])
#
#
#     if uploaded_file:
#         df = pd.read_excel(uploaded_file, skiprows=4)
#         st.markdown("### 🧩 Detected Columns:")
#         st.write(df.columns.tolist())
#
#         df.columns = df.columns.str.strip().str.lower()
#
#         date_col = next((col for col in df.columns if 'date' in col), None)
#         name_col = next((col for col in df.columns if 'name' in col), None)
#         bal_col = next((col for col in df.columns if 'open' in col and 'bal' in col), None)
#
#         if not all([date_col, name_col, bal_col]):
#             st.error("❌ Missing required columns: 'date', 'name', or 'open balance'.")
#             return
#
#         df = df.rename(columns={date_col: "Date", name_col: "Name", bal_col: "Open Balance"})
#         df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
#         df = df.dropna(subset=["Date", "Name"])
#         df["Year"] = df["Date"].dt.year
#         df["Open Balance"] = pd.to_numeric(df["Open Balance"], errors="coerce").fillna(0)
#
#         df_filtered = df[df["Year"].isin([2021, 2022, 2023, 2024])]
#         owing_by_year = {
#             year: df_filtered[(df_filtered["Year"] == year) & (df_filtered["Open Balance"] > 0)]["Name"].unique().tolist()
#             for year in [2021, 2022, 2023, 2024]
#         }
#
#         st.divider()
#
#         # Members Who Owed in All Years
#         owing_all_years = list(set(owing_by_year[2021]) & set(owing_by_year[2022]) & set(owing_by_year[2023]) & set(owing_by_year[2024]))
#         st.subheader("👥 Members Who Owed in 2021, 2022, 2023 and 2024")
#         if owing_all_years:
#             st.write(pd.DataFrame(owing_all_years, columns=["Name"]))
#         else:
#             st.success("✅ No members owed across all three years.")
#
#         unique_members = df_filtered.groupby("Year")["Name"].nunique().to_dict()
#         st.divider()
#
#         st.subheader("📅 Unique Members, Owing Counts & Totals by Year")
#         col1, col2, col3, col4 = st.columns(4)
#         colors = ["#e1f0ff", "#d9f2e6", "#fff4cc", "#fde0dc"]
#
#         for i, (year, col) in enumerate(zip([2021, 2022, 2023, 2024], [col1, col2, col3, col4])):
#             unique_count = unique_members.get(year, 0)
#             owing_names = owing_by_year[year]
#             total_owed = df_filtered[
#                 (df_filtered["Year"] == year) & (df_filtered["Name"].isin(owing_names))
#             ]["Open Balance"].sum()
#
#             col.markdown(f"""
#                 <div style="background-color: {colors[i]}; padding: 20px; border-radius: 10px; box-shadow: 0 1px 4px rgba(0,0,0,0.1);">
#                     <h4 style="color:#0a3d62;">📅 Year: {year}</h4>
#                     <p style="font-size: 16px; color:#333;"><strong>Unique Members:</strong> {unique_count}</p>
#                     <p style="font-size: 16px; color:#333;"><strong>Owing Members:</strong> {len(owing_names)}</p>
#                     <p style="font-size: 16px; color:#333;"><strong>Total Owed:</strong> ${total_owed:,.2f}</p>
#                 </div>
#             """, unsafe_allow_html=True)
#
#         st.divider()
#
#         # Highest Owing Member Each Year
#         st.subheader("💰 Highest Owing Member Each Year")
#         highest_owers = []
#         for year in [2021, 2022, 2023, 2024]:
#             sub_df = df_filtered[(df_filtered["Year"] == year) & (df_filtered["Open Balance"] > 0)]
#             if not sub_df.empty:
#                 top_row = sub_df.loc[sub_df["Open Balance"].idxmax()]
#                 highest_owers.append({"Year": year, "Name": top_row["Name"], "Amount": top_row["Open Balance"]})
#             else:
#                 highest_owers.append({"Year": year, "Name": "—", "Amount": 0.00})
#
#         df_top_owers = pd.DataFrame(highest_owers)
#         df_top_owers["Amount"] = df_top_owers["Amount"].apply(lambda x: f"${x:,.2f}")
#         st.table(df_top_owers.rename(columns={"Year": "📅 Year", "Name": "👤 Name", "Amount": "💸 Amount Owed"}))
#
#         st.divider()
#
#         # Top 5 Owing Each Year
#         with st.expander('💰 Top 5 Highest Owing Members Each Year'):
#             for year in [2021, 2022, 2023, 2024]:
#                 st.markdown(f"### 📅 Year: {year}")
#                 sub_df = df_filtered[(df_filtered["Year"] == year) & (df_filtered["Open Balance"] > 0)]
#                 if not sub_df.empty:
#                     top5 = sub_df[["Name", "Open Balance"]].groupby("Name").sum().sort_values(by="Open Balance", ascending=False).head(5).reset_index()
#                     top5["Open Balance"] = top5["Open Balance"].apply(lambda x: f"${x:,.2f}")
#                     st.table(top5.rename(columns={"Name": "👤 Name", "Open Balance": "💸 Amount Owed"}))
#                 else:
#                     st.info("No members with outstanding balances for this year.")
#
#         # Owing Members List
#         st.subheader("📋 Owing Members and Amounts by Year")
#         def get_name_and_amount(df, year):
#             sub_df = df[(df["Year"] == year) & (df["Open Balance"] > 0)]
#             return [f"{row['Name']} (${row['Open Balance']:,.2f})" for _, row in sub_df.iterrows()]
#
#         data_2021 = get_name_and_amount(df_filtered, 2021)
#         data_2022 = get_name_and_amount(df_filtered, 2022)
#         data_2023 = get_name_and_amount(df_filtered, 2023)
#         data_2024 = get_name_and_amount(df_filtered, 2024)
#         max_rows = max(len(data_2021), len(data_2022), len(data_2023), len(data_2024))
#
#         df_owing_table = pd.DataFrame({
#             "2021": data_2021 + [""] * (max_rows - len(data_2021)),
#             "2022": data_2022 + [""] * (max_rows - len(data_2022)),
#             "2023": data_2023 + [""] * (max_rows - len(data_2023)),
#             "2024": data_2024 + [""] * (max_rows - len(data_2024)),
#         })
#
#         st.dataframe(df_owing_table)
#
#         st.divider()
#
#         # 📈 Members Owing Count Bar Chart
#         st.subheader("📈 Members Owing Per Year")
#         fig, ax = plt.subplots()
#         years = [2021, 2022, 2023, 2024]
#         counts = [len(owing_by_year[year]) for year in years]
#
#         bars = ax.bar(years, counts, color='skyblue')
#         for bar in bars:
#             height = bar.get_height()
#             ax.annotate(f'{int(height)}', xy=(bar.get_x() + bar.get_width() / 2, height),
#                         xytext=(0, 3), textcoords="offset points",
#                         ha='center', va='bottom', fontsize=10)
#
#         ax.set_ylabel("Members Owing")
#         ax.set_xlabel("Year")
#         ax.set_title("Outstanding Dues by Year")
#         st.pyplot(fig)
#
#         st.divider()
#
#         # 📋 Manual Member Management
#         st.subheader("📋 OCA Members – Manage Manually")
#         default_members = [
#             {"Name": "John Doe", "Open Balance": 150.0},
#             {"Name": "Angela Nwosu", "Open Balance": 200.0},
#             {"Name": "Peter Okafor", "Open Balance": 100.0},
#             {"Name": "Chinwe Opara", "Open Balance": 80.0},
#             {"Name": "Samuel Obi", "Open Balance": 50.0},
#         ]
#
#         if "oca_manual_members" not in st.session_state:
#             st.session_state.oca_manual_members = default_members.copy()
#
#         members = st.session_state.oca_manual_members
#         for i, member in enumerate(members):
#             with st.expander(f"👤 {member['Name']} – ${member['Open Balance']:,.2f}"):
#                 col1, col2, col3 = st.columns([4, 2, 1])
#                 new_name = col1.text_input(f"Edit Name #{i}", value=member["Name"], key=f"edit_name_{i}")
#                 new_amount = col2.number_input(f"Edit Amount #{i}", value=member["Open Balance"], step=10.0, key=f"edit_amt_{i}")
#                 delete = col3.button("🗑️ Delete", key=f"delete_{i}")
#
#                 if delete:
#                     st.session_state.oca_manual_members.pop(i)
#                     st.experimental_rerun()
#
#                 member["Name"] = new_name
#                 member["Open Balance"] = new_amount
#
#         st.markdown("### ✅ Updated OCA Members List")
#         df_updated = pd.DataFrame(st.session_state.oca_manual_members)
#         df_updated["Open Balance"] = df_updated["Open Balance"].apply(lambda x: f"${x:,.2f}")
#         st.dataframe(df_updated)
#
# # Run the function
# analyze_membership_debt()



# ===========================================================================================================


#
#
# import streamlit as st
# from deepface import DeepFace
# from PIL import Image
# import tempfile
#
# def emotions_analysis_page():
#     st.title("🧠 Emotion Detection from Face")
#     st.write("Upload a photo and this app will predict the facial emotion!")
#
#     # Upload section
#     uploaded_file = st.file_uploader(
#         "📤 Upload an image (JPG, PNG)",
#         type=["jpg", "jpeg", "png"],
#         key="emotion_uploader"
#     )
#
#     if uploaded_file is not None:
#         # Display uploaded image
#         image = Image.open(uploaded_file).convert("RGB")
#         st.image(image, caption="Uploaded Image", use_column_width=True)
#
#         # Save image to a temporary file
#         with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp:
#             image.save(tmp.name)
#             img_path = tmp.name
#
#         # Analyze emotion
#         with st.spinner("Analyzing emotion..."):
#             try:
#                 result = DeepFace.analyze(img_path=img_path, actions=["emotion"])
#                 dominant_emotion = result[0]['dominant_emotion']
#                 st.success(f"✅ Dominant Emotion: **{dominant_emotion.upper()}**")
#
#                 st.markdown("### 🔍 Emotion Probabilities")
#                 st.json(result[0]['emotion'])
#
#             except Exception as e:
#                 st.error(f"❌ Error analyzing image: {e}")



# =========================================================================================================
#                              OCA ENDS HERE
# ===========================================================================================================



# ============================
#
# import openai
# import streamlit as st
# import os
#
# # 🔐 Load your API key securely
# # openai.api_key = os.getenv("OPENAI_API_KEY")
# openai.api_key = "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"
#
# # ✅ Function to generate email with fallback
# def generate_email(subject, context, tone="professional"):
#     messages = [
#         {"role": "system", "content": "You are an assistant that writes professional emails."},
#         {"role": "user", "content": f"Write a {tone} email about:\nSubject: {subject}\nContext: {context}"}
#     ]
#
#     try:
#         # Try using GPT-4-Turbo first
#         response = openai.ChatCompletion.create(
#             model="gpt-4-turbo",
#             messages=messages,
#             temperature=0.7
#         )
#     except openai.error.InvalidRequestError:
#         # Fall back to GPT-3.5-Turbo if GPT-4 not available
#         response = openai.ChatCompletion.create(
#             model="gpt-3.5-turbo",
#             messages=messages,
#             temperature=0.7
#         )
#
#     return response.choices[0].message['content']
#
#
# # ✅ Streamlit UI
# st.title("✉️ GPT Email Writer")
#
# subject = st.text_input("Email Subject")
# context = st.text_area("Email Context")
# tone = st.selectbox("Select Tone", ["professional", "friendly", "persuasive"])
#
# if st.button("Generate Email"):
#     if not subject or not context:
#         st.warning("Please fill in both Subject and Context.")
#     else:
#         with st.spinner("Generating email..."):
#             result = generate_email(subject, context, tone)
#         st.markdown("### ✅ Generated Email")
#         st.write(result)

#
#

# # ================================================================================ API WRITE END HERE
#
# def sam_data_analysis_page():
#     st.title("📊 SAM Data Analysis")
#     uploaded_file = st.file_uploader("Upload SAM Data File (CSV)", type=['csv'])
#     if uploaded_file:
#         df = pd.read_csv(uploaded_file, encoding='ISO-8859-1')
#         st.write("Data Preview:")
#         st.dataframe(df)
#         # Place your logic here to analyze or summarize `df`
#
# # ======================== PAGE ROUTING ===========================
#
# def main():
#     st.sidebar.title("📂 Navigation")
#     page = st.sidebar.radio("Choose a page:", [
#         "Home",
#         "OCA CONSTITUTION",
#         "OCA MEMBERSHIP",
#         "SAM DATA ANALYSIS"
#         "EMOTIONS"
#     ])
#
#     if page == "Home":
#         home_page()
#     elif page == "OCA CONSTITUTION":
#         oca_constitution_page()
#
#     elif page == "OCA MEMBERSHIP":
#         analyze_membership_debt()
#
#     elif page == "SAM DATA ANALYSIS":
#         sam_data_analysis_page()
#
#
#     elif page == "SAM DATA EMOTIONS":
#         emotions_analysis_page()
# # ========================= LAUNCH APP ============================
#
# if __name__ == "__main__":
#     main()


# =========================================
# 
